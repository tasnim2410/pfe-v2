# backendTry1/app.py
import os
import time
from flask import Flask, request, jsonify
from dotenv import load_dotenv
import pandas as pd
from sqlalchemy import text
# import your scraper & post‐download processor
from scraping_raw_data import EspacenetScraper, process_downloaded_data ,DatabaseManager ,extract_field_keyword_pairs
#from family_members import ensure_columns_exist
from db import db, RawPatent ,ImpactFactor
from family_members import process_dataframe_parallel
from flask import Flask, request, jsonify
import concurrent.futures
from family_members import process_patent_api , PatentsSearch , process_rows
from flask import Flask, request, jsonify
import logging
from outputs.scaler import GlobalScalers

import concurrent.futures
import time
from urllib.parse import quote
import requests
import sqlalchemy
from dotenv import load_dotenv
import threading
import json
from sqlalchemy import text
import uuid
from family_members2 import PatentsSearch , build_espacenet_url
import ast
from keyword_analysis import preprocess_text , extract_keywords ,analyze_topic_evolution
from cleaners import clean_family_members , extract_country_codes  # Import from your module
from family_analysis import get_family_counts_by_country
from research_retrieve import fetch_research_data , process_research_data, store_research_data
from research_retrieve2 import fetch_research_data2 , process_research_data2, store_research_data2
from impact_factor_processor import clean_and_process_data , store_processed_data
from sqlalchemy.exc import SQLAlchemyError
from pandas.errors import ParserError
from research_retrieve2 import fetch_research_data2, process_research_data2, store_research_data2
from research_retrieve3 import (
    fetch_research_data3, process_research_data3, fetch_openalex_works, process_documents,
    clean_issn, clean_doi, renaming_columns, merge_unique_by_doi, store_research_data3,
    filter_by_impact_factor
)
from flask_cors import CORS
from sklearn.feature_extraction.text import TfidfVectorizer
import numpy as np
from flask import jsonify
import math
from collections import Counter ,defaultdict
from keywords_coccurrence_trends import track_cooccurrence_trends , clean_text_remove_stopwords
from growth_rate import patent_growth_summary , compute_patent_growth_from_counts
from keyword_analysis import preprocess_text, analyze_topic_evolution
import logging
from patent_trend_by_field import get_ipc_meaning , normalize_engineering, get_patent_fields
from db import Window, Topic , PatentKeyword ,Divergence ,IPCClassification , ClassifiedPatent , IPCFieldOfStudy , Applicant , ApplicantType , Cost ,PatentCost
from applicant_analysis import get_applicants_df, get_applicant_type_df , get_top_10_applicants,get_innovation_cycle
from market_cost import update_cost_df , calculate_age ,assign_cost, get_market_metrics
import socket
import re
from dotenv import set_key
from top_ipc_classes import get_ipc_codes_meanings
from originality_rate import retrieve_citation_publication_numbers, get_patent_biblio, get_all_citations_ipc
import random
load_dotenv()
import signal
import threading
from werkzeug.serving import make_server
from applicant_analysis import extract_applicant_collaboration_network
from keyword_analysis2 import research_preprocess_text, research_analyze_topic_evolution , build_research_keyword_df
import growth_rate as grmod
from concurrent.futures import ThreadPoolExecutor, as_completed
from threading import Lock
from dotenv import load_dotenv
from urllib.parse import quote
import os, time, requests
from originality_rate import fetch_and_store_originality_data, calculate_originality_rate_from_db, OriginalityData
from datetime import datetime
import uuid
from db import ResearchData3, ResearchWindow, ResearchTopic, ResearchDivergence  # research tables
import subprocess
import json
import os
import time
import threading
import requests
from urllib.parse import quote
from dotenv import load_dotenv
from flask import Flask, jsonify
from tqdm import tqdm   # for progress bar
from sqlalchemy.orm import scoped_session
from db import RawPatent, db , ResearchData3
from family_ops import get_access_token, validate_patent_number, extract_jurisdictions_and_members , load_api_credentials ,fetch_family_data_api , fetch_family_data_scrape
from flask import Flask, request, send_file
from pptx import Presentation
from pptx.util import Inches
import io
import base64
from ops_search import json_to_cql, fetch_to_dataframe ,extract_keyword_pairs
# add near your other imports in app.py
import numpy as np
from tensorflow.keras.models import load_model
from sqlalchemy import text
from flask import request, jsonify
import pandas as pd
import numpy as np
import os, joblib
from datetime import date

# use the SAME feature-engineering as during training
from lstm2 import create_enhanced_features, _y_inv
from growth_rate import compute_patent_growth_from_counts  # past/current growth utils
import joblib
from prophet_single_series import (
    run_for_current_series,
    fetch_current_series,
    PUB_TAIL_TRUNC_DEFAULT,
    PAT_TAIL_TRUNC_DEFAULT,
    TEST_YEARS_DEFAULT,
    HORIZON_DEFAULT,
    MAX_LAG_DEFAULT,
    _require_engine,  
    validate_horizon ,
    compute_patent_growth_from_counts , 
    # to get a DB engine for fetch_current_series
)
from tensorflow.keras.models import load_model
from lstm import build_tech_dataframe, create_sequences_for_tech, _y_inv, fit_scaler_guarded
from sqlalchemy import text
import pandas as pd
from flask import request, jsonify
import uuid
from db import RawPatent, SearchKeyword
import os, joblib
from tensorflow.keras.models import load_model
from sqlalchemy import text
from flask import request, jsonify
import pandas as pd
import numpy as np
import os, joblib
from datetime import date

# === MUST MATCH TRAINING PIPELINE ===
from growth_rate import compute_patent_growth_from_counts  # past/current growth helpers
class ServerThread(threading.Thread):
    def __init__(self, app, port):
        threading.Thread.__init__(self)
        self.srv = make_server('127.0.0.1', port, app)
        self.ctx = app.app_context()
        self.ctx.push()

    def run(self):
        print(f"Running Flask app on http://127.0.0.1:{port}")
        self.srv.serve_forever()

    def shutdown(self):
        self.srv.shutdown()

def create_app():
    
    app = Flask(__name__)
    CORS(app, origins=["http://localhost:3000"])  

    app.config['SQLALCHEMY_DATABASE_URI'] = os.getenv('DATABASE_URL')
    app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
    db.init_app(app)
    # logging.basicConfig(
    # level=logging.INFO,
    # format='%(asctime)s - %(levelname)s - %(message)s'
    # )
    logger = logging.getLogger(__name__)    
    
# Set up logging
    # Database connection
    db_url = os.getenv("DATABASE_URL")
    if db_url is None:
        raise ValueError("DATABASE_URL not found. Please check your .env file.")

# Create the SQLAlchemy engine
    engine = sqlalchemy.create_engine(db_url)

    
    @app.route('/')
    def home():
      return 'app is running!', 200
  
  
      # -------------------- TVP endpoints (health + forecast) --------------------
    # Small, self-contained routes to keep your existing config untouched.

    @app.get("/api/tvp/health")
    def tvp_health():
        """
        Simple liveness check to verify the app is wired and reachable.
        Frontend can ping this before calling the forecast endpoint.
        """
        return jsonify({"ok": True, "service": "tvp", "message": "tvp service is up"}), 200

    @app.post("/api/tvp/forecast")
    def tvp_forecast():
        """
        Runs the R script that emits JSON for plotting.
        Expects your R file to print compact JSON to stdout (cat(toJSON(...))).
        Optional JSON body: { "truncate_last_n": 4, "forecast_h": 10, "lags": [1,2,3] }
        """
        import subprocess, json, os
        payload = request.get_json(silent=True) or {}

        # Resolve the R script path. Keep it configurable without touching other configs.
        # Set TVP_R_SCRIPT_PATH in .env if you keep the file elsewhere.
        backend_dir = os.path.dirname(os.path.abspath(__file__)) if '__file__' in globals() else os.getcwd()
        r_script_path = os.getenv("TVP_R_SCRIPT_PATH", os.path.join(backend_dir, "r", "tvp_var_forecast.R"))
        if not os.path.isfile(r_script_path):
            return jsonify({"ok": False, "error": f"R script not found", "path": r_script_path}), 500

        # Build command (the R script should read commandArgs and output JSON)
        # We pass DB params via env so your R code can pick them up if needed.
        truncate = int(payload.get("truncate_last_n", 4))
        h        = int(payload.get("forecast_h", 10))
        lags     = payload.get("lags", [1, 2, 3])
        if isinstance(lags, list):
            lags_csv = ",".join(str(int(x)) for x in lags)
        else:
            lags_csv = str(lags)

        cmd = [
            "Rscript",
            r_script_path,
            f"--dbname={os.getenv('PG_DB', 'patent_db')}",
            f"--host={os.getenv('PG_HOST', 'localhost')}",
            f"--port={os.getenv('PG_PORT', '5433')}",
            f"--user={os.getenv('PG_USER', 'postgres')}",
            f"--password={os.getenv('PG_PASSWORD', 'tasnim')}",
            f"--truncate_last_n={truncate}",
            f"--forecast_h={h}",
            f"--lags={lags_csv}",
        ]

        try:
            # Run R and capture stdout/stderr
            proc = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                env=os.environ,
                check=False,
                timeout=180
            )
        except FileNotFoundError:
            return jsonify({"ok": False, "error": "Rscript binary not found in PATH"}), 500
        except subprocess.TimeoutExpired:
            return jsonify({"ok": False, "error": "R script timed out"}), 504
        except Exception as e:
            return jsonify({"ok": False, "error": f"Unexpected error: {e}"}), 500

        if proc.returncode != 0:
            # Surface stderr to help you debug R issues quickly
            return jsonify({
                "ok": False,
                "error": "R script failed",
                "return_code": proc.returncode,
                "stderr": proc.stderr.strip()
            }), 502

        stdout = (proc.stdout or "").strip()

        # Extract JSON substring if R prints logs before JSON
        json_start = stdout.find("{")
        json_end   = stdout.rfind("}") + 1
        if json_start != -1 and json_end != -1:
            stdout = stdout[json_start:json_end]

        try:
            data = json.loads(stdout)
        except json.JSONDecodeError as e:
            return jsonify({
                "ok": False,
                "error": "Invalid JSON from R script",
                "detail": str(e),
                "stdout_head": (proc.stdout or "")[:2000]
            }), 500

        # Transform the new R script output structure for backward compatibility
        # and enhanced functionality
        try:
            # Extract models and original data from new structure
            models = data.get("models", {})
            original_data = data.get("original_data", {})
            
            # Create enhanced response structure
            enhanced_data = {
                "models": models,
                "original_data": original_data,
                # Backward compatibility: use first model for legacy endpoints
                "patents": {},
                "publications": {}
            }
            
            # For backward compatibility, use the first available model
            if models:
                first_model_key = list(models.keys())[0]
                first_model = models[first_model_key]
                
                # Combine historical and forecast data for backward compatibility
                if "historical" in first_model and "forecast" in first_model:
                    hist_pat = first_model["historical"]
                    fore_pat = first_model["forecast"]
                    hist_pub = first_model["historical"]
                    fore_pub = first_model["forecast"]
                    
                    enhanced_data["patents"] = {
                        "years": hist_pat.get("years", []) + fore_pat.get("years", []),
                        "values": hist_pat.get("patent_count", []) + fore_pat.get("patent_count", [])
                    }
                    enhanced_data["publications"] = {
                        "years": hist_pub.get("years", []) + fore_pub.get("years", []),
                        "values": hist_pub.get("pub_count", []) + fore_pub.get("pub_count", [])
                    }
            
            return jsonify({"ok": True, "data": enhanced_data}), 200
            
        except Exception as transform_error:
            # If transformation fails, return original data
            return jsonify({"ok": True, "data": data, "transform_warning": str(transform_error)}), 200



  
    @app.post("/api/arimax/forecast")
    def arimax_quadratic_forecast():
        """
        Shells out to the R script and returns its JSON.
        Accepts optional body:
        { "mode":"test","split_year":2015 }
        { "mode":"future","horizon":5,"pub_future":{"years":[...],"values":[...]}, "pub_future_strategy":"linear" }
        """
        import subprocess, json, os
        payload = request.get_json(silent=True) or {}

        mode = str(payload.get("mode", "test")).lower()
        split_year = int(payload.get("split_year", 2015))
        horizon = int(payload.get("horizon", 5))
        pub_future = payload.get("pub_future") or {}
        pub_strategy = str(payload.get("pub_future_strategy", "linear")).lower()

        backend_dir = os.path.dirname(os.path.abspath(__file__)) if '__file__' in globals() else os.getcwd()
        r_script_path = os.getenv("ARIMAX_R_SCRIPT_PATH", os.path.join(backend_dir, "r", "arimax_quadratic.R"))
        if not os.path.isfile(r_script_path):
            return jsonify({"ok": False, "error": "R script not found", "path": r_script_path}), 500

        fut_years_csv = fut_vals_csv = ""
        if isinstance(pub_future, dict):
            yrs = pub_future.get("years") or []
            vals = pub_future.get("values") or []
            if isinstance(yrs, list) and isinstance(vals, list) and len(yrs) == len(vals) and len(yrs) > 0:
                fut_years_csv = ",".join(str(int(y)) for y in yrs)
                fut_vals_csv  = ",".join(str(float(v)) for v in vals)

        cmd = [
            "Rscript", r_script_path,
            f"--dbname={os.getenv('PG_DB', 'patent_db')}",
            f"--host={os.getenv('PG_HOST', 'localhost')}",
            f"--port={os.getenv('PG_PORT', '5433')}",
            f"--user={os.getenv('PG_USER', 'postgres')}",
            f"--password={os.getenv('PG_PASSWORD', 'tasnim')}",
            "--sslmode=" + os.getenv("PG_SSLMODE", "prefer"),
            f"--mode={mode}",
            f"--split_year={split_year}",
            f"--horizon={horizon}",
            f"--pub_future_strategy={pub_strategy}",
            "--debug=" + os.getenv("ARIMAX_DEBUG", "0"),
        ]
        if fut_years_csv and fut_vals_csv:
            cmd += [f"--pub_future_years={fut_years_csv}", f"--pub_future_values={fut_vals_csv}"]

        try:
            proc = subprocess.run(cmd, capture_output=True, text=True, env=os.environ, check=False, timeout=180)
        except FileNotFoundError:
            return jsonify({"ok": False, "error": "Rscript not found in PATH"}), 500
        except subprocess.TimeoutExpired:
            return jsonify({"ok": False, "error": "R script timed out"}), 504
        except Exception as e:
            return jsonify({"ok": False, "error": f"Unexpected error: {e}"}), 500

        if proc.returncode != 0:
            return jsonify({"ok": False, "error": "R script failed", "return_code": proc.returncode,
                            "stderr": (proc.stderr or "").strip()}), 502

        stdout = (proc.stdout or "").strip()
        json_start = stdout.find("{"); json_end = stdout.rfind("}") + 1
        if json_start != -1 and json_end != -1:
            stdout = stdout[json_start:json_end]

        try:
            data = json.loads(stdout)
        except json.JSONDecodeError as e:
            return jsonify({"ok": False, "error": "Invalid JSON from R", "detail": str(e),
                            "stdout_head": (proc.stdout or "")[:2000], "stderr": (proc.stderr or "")[:500]}), 500

        return jsonify({"ok": True, "data": data, "stderr": (proc.stderr or "").strip()}), 200

  
  
  
  
  


    @app.route("/api/search_ops", methods=["POST"])
    def search_patents_ops():
        """
        Handles advanced patent search via EPO OPS API.
        Accepts JSON body: { "query": { … }, "max_results": <int> }
        Returns records augmented with computed years and abstract.
        """
        data        = request.get_json(silent=True) or {}
        q_input     = data.get("query")
        max_results = int(data.get("max_results", 500))

        if not q_input:
            return jsonify({"error": "Request must contain 'query' key."}), 400

        # Build and execute CQL
        start_time = time.perf_counter()
        try:
            cql = json_to_cql({"query": q_input})
            app.logger.debug(f"Generated CQL: {cql}")
            df, total_cnt = fetch_to_dataframe(cql, max_records=max_results)
        except Exception as e:
            app.logger.error(f"OPS request failed: {e}")
            return jsonify({"error": f"OPS request failed: {e}"}), 502

        elapsed = time.perf_counter() - start_time
        app.logger.info(f"OPS fetch took {elapsed:.2f}s")

        if df.empty:
            app.logger.info("No results found")
            return jsonify({"error": "no results"}), 404

        # Debug: log DataFrame columns to verify field names
        app.logger.debug(f"DataFrame columns: {df.columns.tolist()}")

        # Prepare unique search ID
        search_id = str(uuid.uuid4())

        # Convert dataframe to list of dicts and augment
        records = df.where(pd.notnull(df), None).to_dict(orient="records")
        for rec in records:
            # Debug: log record keys for mapping
            app.logger.debug(f"Record keys: {list(rec.keys())}")

            # First filing year from 'Publication date'
            pub = rec.get("Publication date")
            try:
                pub_dt = datetime.strptime(pub, "%Y%m%d") if pub else None
                rec["first_filing_year"] = pub_dt.year if pub_dt else None
            except Exception:
                rec["first_filing_year"] = None

            # Earliest priority year: try multiple possible keys
            prio_raw = rec.get("Earliest priority") or rec.get("Earliest priority date")
            try:
                prio_dt = datetime.strptime(prio_raw, "%Y-%m-%d") if prio_raw else None
                rec["earliest_priority_year"] = prio_dt.year if prio_dt else None
            except Exception:
                rec["earliest_priority_year"] = None

            # Abstract: try uppercase/lowercase
            abstract = rec.get("Abstract") or rec.get("abstract") or rec.get("abstractText")
            rec["abstract"] = abstract

        # Store raw patents including new fields
        mappings = []
        for r in records:
            # parse publication_date again for DB
            pub_dt = None
            try:
                pub_dt = datetime.strptime(r.get("Publication date"), "%Y%m%d")
            except Exception:
                pass

            mappings.append({
                "title":                    r.get("Title"),
                "inventors":                r.get("Inventors"),
                "applicants":               r.get("Applicants"),
                "publication_number":       r.get("Publication number"),
                "publication_date":         pub_dt,
                "ipc":                      r.get("IPC"),
                "cpc":                      r.get("CPC"),
                "applicant_country":        r.get("app_country") or r.get("Applicant country"),
                "family_id":            r.get("Family number"),
                "first_publication_number": r.get("Publication number"),
                "first_publication_country": r.get("Publication number", "")[:2],
                "first_filing_year":        r.get("first_filing_year"),
                "earliest_priority_year":   r.get("earliest_priority_year"),
                "abstract":                 r.get("abstract"),
            })

        db.session.query(RawPatent).delete()
        db.session.bulk_insert_mappings(RawPatent, mappings)

        # Extract & store keywords
        group1 = q_input.get("group1") if isinstance(q_input, dict) else None
        for field, word in extract_keyword_pairs(group1) if group1 else []:
            db.session.add(SearchKeyword(
                search_id     = search_id,
                field         = field,
                keyword       = word,
                total_results = total_cnt
            ))
        db.session.commit()

        # Return augmented records
        return jsonify({
            "search_id":     search_id,
            "rows":          len(records),
            "total_results": total_cnt,
            "data":          records
        }), 200






  
    @app.route('/api/last_search_keywords', methods=['GET'])
    def get_last_search_keywords():
        try:
            with engine.connect() as connection: 
                result = connection.execute(
                    text("SELECT search_id FROM search_keywords WHERE id =(SELECT MAX(id) FROM search_keywords)")
                )
                row = result.fetchone()
                if not row: 
                    return jsonify({"message":"No search keywords found"}), 404
                search_id = row[0]

                result = connection.execute(
                    text("SELECT keyword , field FROM search_keywords WHERE search_id = :search_id"),
                    {"search_id": search_id}
                )
                rows = result.mappings().all()
                search_query = {row["keyword"]:row["field"] for row in rows}
                return jsonify(search_query), 200
        except Exception as e:
            return jsonify({"error": f"Failed to fetch last search keywords: {str(e)}"}), 500

    #search keywords with field mapping
     
    # @app.route('/api/search', methods=['POST'])
    # def search_patents():
    #     # 1) Parse complex query parameters from JSON body
    #     data = request.get_json()
    #     keywords_list = data.get('keywords', []) if data else []
    #     if not keywords_list:
    #         return jsonify({"error": "Use JSON body with 'keywords' list: [{\"field\": \"field_name\", \"keyword\": \"keyword\"}]"}), 400
    
    #     search_id = str(uuid.uuid4())
    
    #     # Build search_map from the list of objects
    #     search_map = {}
    #     for item in keywords_list:
    #         field = item.get('field')
    #         keyword = item.get('keyword')
    #         if not field or not keyword:
    #             return jsonify({"error": "Each keyword item must have 'field' and 'keyword'"}), 400
    #         search_map[keyword.strip()] = field.strip().lower()  # Normalize field name

    #     # 2) Build scraper with the parsed query
    #     db_manager = DatabaseManager()
    #     scraper = EspacenetScraper(
    #         search_map,
    #         headless=True,
    #         options_args=[
    #             "--disable-blink-features=AutomationControlled",
    #             "--no-sandbox",
    #             "--disable-dev-shm-usage",
    #             "--remote-debugging-port=9222",
    #             "--user-agent=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/90.0.4430.212 Safari/537.36"
    #         ]
    #     )
    #     try:
    #         # Load the Espacenet page
    #         if not scraper.get_page_html():
    #             return jsonify({"error": "Failed to load Espacenet page"}), 500

    #         # Download the CSV file
    #         if not scraper.download_csv(max_results=500):
    #             return jsonify({"error": "Failed to download CSV"}), 500

    #         # Wait to ensure file download completes
    #         time.sleep(10)

    #         # Process the downloaded CSV
    #         df = process_downloaded_data(os.path.expanduser("~/Downloads"))
    #         if df is None or df.empty:
    #             return jsonify({"error": "Couldn’t parse downloaded CSV or DataFrame is empty"}), 500

    #         print(f"DataFrame contains {len(df)} rows")

    #         # Store data in the database
    #         if not db_manager.store_patents(df):
    #             return jsonify({"error": "Failed to store data in database"}), 500

    #         print("Data stored successfully")

    #     finally:
    #         if scraper:
    #             scraper.close()

    #     # 3) Insert search keywords into the database
    #     try:
    #         with db_manager.engine.connect() as connection:
    #             keyword_params = [
    #                 {"search_id": search_id, "field": field, "keyword": keyword}
    #                 for keyword, field in search_map.items()
    #             ]
    #             connection.execute(
    #                 text("INSERT INTO search_keywords (search_id, field, keyword) VALUES (:search_id, :field, :keyword)"),
    #                 keyword_params
    #             )
    #             connection.commit()
    #             print("Keyword data inserted successfully")
    #     except Exception as e:
    #         print(f"Error inserting keywords into database: {e}")
    #         return jsonify({"error": f"Failed to insert search keywords: {str(e)}"}), 500

    #     df = df.where(pd.notnull(df), None)

    #     for dtcol in df.select_dtypes(include=['datetime64[ns]']).columns:
    #         df[dtcol] = df[dtcol].apply(lambda x: x.isoformat() if x is not None else None)
    
    #     # 4) Return the response
    #     response_data = {
    #         "search_id": search_id,
    #         "results": df.to_dict(orient='records')
    #     }

    #     return jsonify(response_data), 200
        def safe_number(val):
            return val if isinstance(val, (int, float)) and not math.isnan(val) else None

    
    # @app.route('/api/search', methods=['POST'])
    # def search_patents():
    #     data = request.get_json()
    #     keywords_list = data.get('keywords', [])
    #     operators = data.get('operators', [])
    #     if not keywords_list or len(keywords_list) < 2:
    #         return jsonify({"error": "Provide at least two keywords."}), 400
    #     if len(operators) != len(keywords_list) - 1:
    #         return jsonify({"error": "Operators count must be one less than keywords count."}), 400

    #     search_id = str(uuid.uuid4())
    #     search_pairs = []
    #     for item in keywords_list:
    #         field = item.get('field')
    #         keyword = item.get('keyword')
    #         if not field or not keyword:
    #             return jsonify({"error": "Each keyword item must have 'field' and 'keyword'"}), 400
    #         search_pairs.append((field.strip().lower(), keyword.strip()))

    #     db_manager = DatabaseManager()
    #     scraper = EspacenetScraper(
    #         search_pairs,
    #         operators=operators,
    #         headless=True,
    #         options_args=[
    #             "--disable-blink-features=AutomationControlled",
    #             "--no-sandbox",
    #             "--disable-dev-shm-usage",
    #             "--remote-debugging-port=9222",
    #             "--user-agent=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/90.0.4430.212 Safari/537.36"
    #         ]
    #     )
    #     # ...rest of your code...
    #     try:
    #         # Load the Espacenet page
    #         if not scraper.get_page_html():
    #             return jsonify({"error": "Failed to load Espacenet page"}), 500

    #         # Download the CSV file
    #         if not scraper.download_csv(max_results=500):
    #             return jsonify({"error": "Failed to download CSV"}), 500

    #         # Wait to ensure file download completes
    #         time.sleep(10)

    #         # Process the downloaded CSV
    #         df = process_downloaded_data(os.path.expanduser("~/Downloads"))
    #         if df is None or df.empty:
    #             return jsonify({"error": "Couldn’t parse downloaded CSV or DataFrame is empty"}), 500

    #         print(f"DataFrame contains {len(df)} rows")

    #         # Store data in the database
    #         if not db_manager.store_patents(df):
    #             return jsonify({"error": "Failed to store data in database"}), 500

    #         print("Data stored successfully")

    #     finally:
    #         if scraper:
    #             scraper.close()

    #     # Create search_map from search_pairs for database storage
    #     search_map = {keyword: field for field, keyword in search_pairs}

    #     # 3) Insert search keywords into the database
    #     try:
    #         with db_manager.engine.connect() as connection:
    #             keyword_params = [
    #                 {"search_id": search_id, "field": field, "keyword": keyword}
    #                 for keyword, field in search_map.items()
    #             ]
    #             connection.execute(
    #                 text("INSERT INTO search_keywords (search_id, field, keyword) VALUES (:search_id, :field, :keyword)"),
    #                 keyword_params
    #             )
    #             connection.commit()
    #             print("Keyword data inserted successfully")
    #     except Exception as e:
    #         print(f"Error inserting keywords into database: {e}")
    #         return jsonify({"error": f"Failed to insert search keywords: {str(e)}"}), 500

    #     df = df.replace({np.nan: None})
        
    #     for dtcol in df.select_dtypes(include=['datetime64[ns]']).columns:
    #         df[dtcol] = df[dtcol].apply(lambda x: x.isoformat() if x is not None else None)
    
    #     # 4) Return the response
    #     response_data = {
    #         "search_id": search_id,
    #         "results": df.to_dict(orient='records')  
    #     }

    #     return jsonify(response_data), 200
    
    
    
    
    
    
    
    @app.route('/api/search', methods=['POST'])
    def search_patents():
            query_input = request.get_json()

            db_manager = DatabaseManager()
            scraper = EspacenetScraper(
                search_keywords=None,
                headless=True,
            )
           
            try:
                # Load the Espacenet page
                if not scraper.get_page_html(query_input):
                    return jsonify({"error": "Failed to load Espacenet page"}), 500

                # Download the CSV file
                if not scraper.download_csv(max_results=500):
                    return jsonify({"error": "Failed to download CSV"}), 500

                # Wait to ensure file download completes
                time.sleep(10)

                # Process the downloaded CSV
                df = process_downloaded_data(os.path.expanduser("~/Downloads"))
                if df is None or df.empty:
                    return jsonify({"error": "Couldn’t parse downloaded CSV or DataFrame is empty"}), 500

                print(f"DataFrame contains {len(df)} rows")

                # Store data in the database
                if not db_manager.store_patents(df):
                    return jsonify({"error": "Failed to store data in database"}), 500

                print("Data stored successfully")

            finally:
                if scraper:
                    scraper.close()

            root_group = query_input.get("query", {}).get("group1")
            if not root_group:
                return jsonify({"error": "Query must contain a group labeled 'group1' under 'query'"}), 400
            search_pairs = extract_field_keyword_pairs(root_group)
            search_id = str(uuid.uuid4())

            # 3) Insert search keywords into the database
            try:
                with db_manager.engine.connect() as connection:
                    keyword_params = [
                        {"search_id": search_id, "field": field, "keyword": keyword}
                        for field, keyword in search_pairs
                    ]
                    connection.execute(
                    text("""
                     INSERT INTO search_keywords (search_id, field, keyword)
                 VALUES (:search_id, :field, :keyword)
                """),
                    keyword_params
                )
                    connection.commit()
                    print("Keyword data inserted successfully")
            except Exception as e:
                print(f"Error inserting keywords into database: {e}")
                return jsonify({"error": f"Failed to insert search keywords: {str(e)}"}), 500

            df = df.replace({np.nan: None})
        
            for dtcol in df.select_dtypes(include=['datetime64[ns]']).columns:
                df[dtcol] = df[dtcol].apply(lambda x: x.isoformat() if x is not None else None)
    
            # 4) Return the response
            response_data = {
                "search_id": search_id,
                "results": df.to_dict(orient='records')  
            }

            return jsonify(response_data), 200
    

    
    
    
    

    def fetch_last_search_keywords_from_db():
        """
        Returns a dict mapping keyword→field for the most
        recent search_id, or raises ValueError if none found.
        """
        with engine.connect() as conn:
            row = conn.execute(
                text("SELECT search_id FROM search_keywords ORDER BY id DESC LIMIT 1")
            ).fetchone()
            if not row:
                raise ValueError("No search keywords found")

            search_id = row[0]
            mappings = conn.execute(
                text("SELECT keyword, field FROM search_keywords WHERE search_id = :sid"),
                {"sid": search_id}
            ).mappings().all()

        # build and return the dict
            return {m["keyword"]: m["field"] for m in mappings}
        

    @app.route('/api/family_members/scraping', methods=['POST'])
    def get_family_members():
    # Fetch data from database
        field_mapping = {
            "title": "ti",
            "abstract": "ab",
            "claims": "cl",
            "title,abstract or claims": "ctxt",
            "all text fields": "ftxt",
            "title or abstract": "ta",
            "description": "desc",
            "all text fields or names": "nftxt",
            "title , abstract or names": "ntxt"
        }
        keywords = fetch_last_search_keywords_from_db()
    
        # Load existing data
        query = 'SELECT *, "No" as id FROM raw_patents'
        df = pd.read_sql(query, engine)
    
        scraper = PatentsSearch(headless=False)
    
        try:
            # Process each patent record
            for index, row in df.iterrows():
                print(f"Processing patent {index+1}/{len(df)}")
            
            # 1. Scrape family members
                url = build_espacenet_url(row, keywords, field_mapping)
                html = scraper.get_page_html(url)
            
                if not html:
                    print(f"Skipping {row['first_publication_number']} - failed to retrieve page")
                    df.at[index, 'family_members'] = []
                    df.at[index, 'family_jurisdictions'] = []
                    continue
                
            # 2. Get raw family members from HTML
                raw_members = scraper.parse_html(html)
            
            # 3. Clean the scraped data
                cleaned_members = clean_family_members(raw_members)
            
            # 4. Extract country codes from CLEANED members
                jurisdictions = extract_country_codes(cleaned_members)
            
            # 5. Update both columns
                df.at[index, 'family_members'] = cleaned_members
                df.at[index, 'family_jurisdictions'] = jurisdictions
            
                print(f"Updated {row['first_publication_number']} with {len(cleaned_members)} members")

            # Prepare database updates
            updates = []
            for _, row in df.iterrows():
                updates.append({
                    'id': row['id'],
                    'members': row['family_members'],
                    'jurisdictions': row['family_jurisdictions']
                })

            # Update database in a single transaction
            with db.engine.begin() as connection:  # Automatically commits/rolls back
                connection.execute(
                    text("""
                        UPDATE raw_patents 
                        SET family_members = :members,
                            family_jurisdictions = :jurisdictions
                        WHERE "No" = :id
                    """),
                    updates
                )

            # Prepare response statistics
            total_members = sum(len(m) for m in df['family_members'])
            unique_countries = list({cc for codes in df['family_jurisdictions'] for cc in codes})
        
            return jsonify({
                "success": True,
                "updated_records": len(df),
                "total_family_members": total_members,
                "unique_jurisdictions": sorted(unique_countries),
                "sample_entry": {
                    "publication_number": df.iloc[0]['first_publication_number'],
                    "members": df.iloc[0]['family_members'],
                    "jurisdictions": df.iloc[0]['family_jurisdictions']
                }
            })

        except Exception as e:
            return jsonify({
                "success": False,
                "error": str(e),
                "failed_at": f"Failed at record {index}: {row['first_publication_number']}" if 'index' in locals() else "Initialization"
            }), 500

        finally:
            if 'scraper' in locals():
                scraper.close()

    @app.route('/api/research', methods=['POST'])
    def store_research():
        query = request.json.get('query')
        if not query:
            return jsonify({"error": "Query parameter is required"}), 400
    
        papers = fetch_research_data(query)
        if not papers:
            return jsonify({"error": "Failed to fetch data from API"}), 500
    
        df = process_research_data(papers)
        try:
            store_research_data(df)
            return jsonify({"message": "Data stored successfully"}), 200
        except Exception as e:
            return jsonify({"error": str(e)}), 500
        
        
        
    @app.route('/api/research_retrieve_and_filtering', methods=['POST'])
    def store_research2():
        query = request.json.get('query')
        papers = fetch_research_data2(query)
        impact_factors = ImpactFactor.query.all()
    
        if not papers:
            return jsonify({"error": "No papers found"}), 404
    
        df = process_research_data2(papers , impact_factors)
    
        try:
            store_research_data2(df)
            return jsonify({
                "message": f"Stored {len(df)} papers",
                "stats": {
                    
                    "top_journal": df['publication_venue_name'].mode()[0]
                }
            }), 200
        except Exception as e:
            return jsonify({"error": str(e)}), 500
        
        
    """enpoint to fetch and store research data that's supposed to show the number of papers before and after filtering but not working yet"""
    # @app.route('/api/research_retrieve', methods=['POST'])
    # def store_research2():
    #     query = request.json.get('query')
    #     papers = fetch_research_data(query)

    #     if not papers:
    #         return jsonify({"error": "No papers found"}), 404

    # # Must match the two‐value return from process_research_data:
    #     result = process_research_data(papers)
    #     print(f"Received result type: {type(result)}, length: {len(result) if hasattr(result, '__len__') else 'N/A'}")
    #     if len(result) != 2:
    #         return jsonify({"error": "Unexpected data format from processing"}), 500
    #     df, stats = result
    #     if df.empty:
    #         return jsonify({"error": "No papers survived processing filters"}), 40

    #     try:
    #         store_research_data(df)
    #         return jsonify({
    #             "message": f"Stored {len(df)} papers",
    #             "stats": {
    #                 "before_clean": stats['before_clean'],
    #                 "after_clean": stats['after_clean'],
    #                 "after_filter": stats['after_filter'],
    #                 "top_journal": df['publication_venue_name'].mode()[0] if len(df) else None
    #             }
    #         }), 200

    #     except Exception as e:
    #         return jsonify({"error": str(e)}), 500





        
    # @app.route('/process_and_store', methods=['GET'])
    # def process_and_store():
    #     """
    #     Endpoint to check if store_processed_data works.
    #     Requires ?confirm=true query parameter to proceed.
    #     Returns JSON response indicating success or failure.
    #     """
    #     confirm = request.args.get('confirm', 'false').lower() == 'true'
    #     if not confirm:
    #         return jsonify({"message": "Operation not confirmed. Add ?confirm=true to proceed."}), 400
    
    #     try:
    #         logging.info("Starting data processing and storage.")
    #         message = store_processed_data()
    #         logging.info("Data processed and stored successfully.")
    #         return jsonify({"message": message}), 200
    #     except ValueError as e:
    #         logging.error(f"Configuration error: {str(e)}")
    #         return jsonify({"error": f"Configuration error: {str(e)}"}), 400
    #     except FileNotFoundError as e:
    #         logging.error(f"File not found: {str(e)}")
    #         return jsonify({"error": f"File not found: {str(e)}"}), 404
    #     except ParserError as e:
    #         logging.error(f"Error parsing file: {str(e)}")
    #         return jsonify({"error": f"Error parsing file: {str(e)}"}), 400
    #     except SQLAlchemyError as e:
    #         logging.error(f"Database error: {str(e)}")
    #         return jsonify({"error": f"Database error: {str(e)}"}), 500
    #     except Exception as e:
    #         logging.error(f"Unexpected error: {str(e)}")
    #         return jsonify({"error": f"Unexpected error: {str(e)}"}), 500




    # @app.route('/fetch_research_data', methods=['POST'])
    # def fetch_research_data():
    #     # Get query from the request
    #     data = request.get_json()
    #     query = data.get('query')

    #     if not query:
    #         return jsonify({'error': 'Query is required'}), 400

    #     # Load impact factors from the database for processing
    #     impact_factors_df = ImpactFactor.query.all()

    #     # Fetch and process Semantic Scholar data
    #     papers = fetch_research_data3(query)
    #     if not papers:
    #         return jsonify({'error': 'No papers found from Semantic Scholar or API error occurred'}), 500
    #     sem_df = process_research_data3(papers, impact_factors_df)

    #     # Add DOI cleaning for Semantic Scholar DataFrame
    #     sem_df['doi_clean'] = sem_df['DOI'].apply(clean_doi)

    #     # Fetch and process OpenAlex data
    #     openalex_works = fetch_openalex_works(max_docs=300, per_page=200)
    #     if not openalex_works:
    #         return jsonify({'error': 'No works found from OpenAlex or API error occurred'}), 500
    
    #     # Use the new process_documents function with impact_factors_df
    #     openalex_df = process_documents(openalex_works, journals_df=impact_factors_df)

    #     # Clean DOI for OpenAlex DataFrame (ISSN cleaning is now handled in process_documents)
    #     openalex_df['doi_clean'] = openalex_df['doi'].apply(clean_doi)

    #     # Rename columns to ensure consistency before merging
    #     sem_df, openalex_df = renaming_columns(sem_df, openalex_df)

    #     # Merge the two DataFrames using merge_unique_by_doi
    #     final_df = merge_unique_by_doi(sem_df, openalex_df)

    #     # Store the merged DataFrame
    #     store_research_data3(final_df)

    #     # Return a success response
    #     return jsonify({
    #         'message': 'Research data fetched, merged, and stored successfully',
    #         'papers_processed': len(final_df),
    #         'semantic_scholar_papers': len(sem_df),
    #         'openalex_papers': len(openalex_df)
    #     }), 200
        
        
        
    @app.route('/api/scientific_search_merge', methods=['POST'])
    def scientific_search_merge():
        """
        Accepts: { "query": <search_string> }
        Searches both Semantic Scholar and OpenAlex with the query,
        Filters and merges both datasets, stores merged result in research_data3 table.
        Returns a summary JSON.
        """
        try:
            # Parse JSON input
            data = request.get_json()
            query = data.get("query")
            if not query:
                return jsonify({"error": "Query is required"}), 400

            # Load impact factors from DB for enrichment
            impact_factors_df = db.session.query(ImpactFactor).all()

            # Fetch Semantic Scholar papers (S2)
            s2_papers = fetch_research_data3(query)
            print("Semantic Scholar: records before filtering:", len(s2_papers))
            if not s2_papers:
                return jsonify({'error': 'No papers found from Semantic Scholar or API error occurred'}), 500
            s2_df = process_research_data3(s2_papers, impact_factors_df)
            print("Semantic Scholar: records after processing:", len(s2_df))
            # Clean DOI in S2
            s2_df['doi_clean'] = s2_df['DOI'].apply(clean_doi)

            # Fetch OpenAlex works
            openalex_works = fetch_openalex_works(query, max_docs=300, per_page=200)
            print("OpenAlex: records before filtering:", len(openalex_works))

            if not openalex_works:
                return jsonify({'error': 'No works found from OpenAlex or API error occurred'}), 500
            openalex_df = process_documents(openalex_works, journals_df=impact_factors_df)
            openalex_df = filter_by_impact_factor(openalex_df, impact_factors_df)
            print("OpenAlex: records after filtering:", len(openalex_df))
            openalex_df['doi_clean'] = openalex_df['doi'].apply(clean_doi)

            # Standardize columns and merge by DOI
            s2_df, openalex_df = renaming_columns(s2_df, openalex_df)
            final_df = merge_unique_by_doi(s2_df, openalex_df)

            # Store in DB (research_data3 table)
            print("About to store", len(final_df), "records")
            print(final_df.head())

            store_research_data3(final_df)

            return jsonify({
                'message': 'Scientific data fetched, merged, and stored successfully',
                'papers_processed': len(final_df),
                'semantic_scholar_papers': len(s2_df),
                'openalex_papers': len(openalex_df)
            }), 200

        except Exception as e:
            import traceback
            print(traceback.format_exc())
            return jsonify({"error": str(e)}), 500

        
    @app.route('/api/research_field_trends', methods=['POST'])
    def research_field_trends():
        """
        Returns count per year for selected fields of study from research_data3.
        Expects JSON: { "fields": ["field1", "field2", ...] }
        """
        try:
            data = request.get_json()
            selected_fields = data.get("fields")
            if not selected_fields or not isinstance(selected_fields, list):
                return jsonify({"error": "Request must include a list of 'fields'"}), 400

            # Query all records with year and fields_of_study
            results = db.session.query(ResearchData3.year, ResearchData3.fields_of_study).filter(
                ResearchData3.year.isnot(None),
                ResearchData3.fields_of_study.isnot(None)
            ).all()

            # Count for each field and year
            counts = {}
            for year, field_list in results:
                if not isinstance(field_list, list):
                    # If stored as string, try to parse (Postgres can return list or stringified list)
                    import ast
                    try:
                        field_list = ast.literal_eval(field_list)
                    except Exception:
                        continue
                for field in field_list:
                    if field in selected_fields:
                        key = (year, field)
                        counts[key] = counts.get(key, 0) + 1

            # Prepare for frontend
            data_out = [
                {"year": year, "field": field, "count": count}
                for (year, field), count in counts.items()
            ]
            # Optionally, sort by year and field
            data_out.sort(key=lambda x: (x["field"], x["year"]))

            return jsonify(data_out), 200

        except Exception as e:
            import traceback
            print(traceback.format_exc())
            return jsonify({"error": str(e)}), 500
      
        
        
        
        
        
        
        
        
        
        
    @app.route('/top_keyword', methods=['GET'])
    def keyword_analysis():
        
        
        try:
            # Fetch patent data from the database
            query = 'SELECT * FROM raw_patents'
            logger.info("Fetching patent data from raw_patents table")
            df = pd.read_sql(query, engine)
        
            if df.empty:
                logger.warning("No patent data found in the database")
                return jsonify({"error": "No patent data found in the database"}), 404
        
            # Preprocess the Titre column
            logger.info("Preprocessing patent titles")
            df['title'] = df['Title'].apply(preprocess_text)
        
            # Extract texts, dropping any nulls
            texts = df['title'].dropna().tolist()
            if not texts:
                logger.warning("No valid texts after preprocessing")
                return jsonify({"error": "No valid texts available after preprocessing"}), 400
        
            # Extract top keywords
            logger.info(f"Extracting top keywords from {len(texts)} texts")
            top_keywords = extract_keywords(texts)
        
            if not top_keywords:
                logger.warning("No keywords extracted (possibly due to insufficient documents)")
                return jsonify({"warning": "No keywords extracted, try adjusting min_df or adding more data"}), 200
        
            # Format keywords for JSON response
            keywords_response = [{"keyword": keyword, "score": float(score)} for keyword, score in top_keywords]
        
            logger.info("Successfully extracted keywords")
            return jsonify({
                "status": "success",
                "keywords": keywords_response,
                "count": len(keywords_response)
            }), 200
    
        except sqlalchemy.exc.SQLAlchemyError as e:
            logger.error(f"Database error: {str(e)}")
            return jsonify({"error": f"Database error: {str(e)}"}), 500
        except Exception as e:
            logger.error(f"Unexpected error: {str(e)}")
            return jsonify({"error": f"Unexpected error: {str(e)}"}), 500
        
        
                    
    @app.route('/api/processed_texts', methods=['GET'])
    def get_processed_texts_from_db():
            """
            1) Fetch the 'Title' column from raw_patents.
            2) Preprocess each title with preprocess_text(...).
            3) Return JSON: { "processed": [ "...", "...", ... ] }.
            """
            try:
                # 1) Read only the Title column from raw_patents
                query = 'SELECT "Title" FROM raw_patents'
                df = pd.read_sql(query, engine)

                if df.empty:
                    return jsonify({"error": "No records found in raw_patents"}), 404

                # 2) Apply preprocess_text to every Title, skipping NaN
                processed_list = []
                for raw_title in df['Title'].fillna(""):
                    proc = preprocess_text(raw_title, use_stemming=True)
                    if proc: 
                        processed_list.append(proc)

                if not processed_list:
                    return jsonify({"error": "No valid text to process"}), 400

                # 3) Return JSON array
                return jsonify({"processed": processed_list}), 200

            except sqlalchemy.exc.SQLAlchemyError as e:
                # Database‐level errors
                logger.error(f"Database error in /api/processed_texts: {e}")
                return jsonify({"error": f"Database error: {str(e)}"}), 500

            except Exception as e:
                # Any other unexpected errors
                logger.error(f"Unexpected error in /api/processed_texts: {e}")
                return jsonify({"error": f"Unexpected error: {str(e)}"}), 500
            
            
            
    @app.route('/api/patents/first_filing_years', methods=['GET'])
    def get_first_filing_years():
    
        try:
            # Simple query using the existing first_filing_year column
            query = text('''
                SELECT first_filing_year AS year, COUNT(*) AS count
                FROM raw_patents
                WHERE first_filing_year IS NOT NULL
                GROUP BY first_filing_year
                ORDER BY year
            ''')

            with engine.connect() as conn:
                result = conn.execute(query)
                rows = result.fetchall()

            if not rows:
                return jsonify({"message": "No patent data available"}), 404

            # Prepare data for chart
            years = []
            counts = []
            for row in rows:
                years.append(int(row[0]))
                counts.append(row[1])

            return jsonify({
                "labels": years,
                "datasets": [{
                    "label": "Number of Patents",
                    "data": counts,
                    "borderColor": "rgb(75, 192, 192)",
                    "tension": 0.1
                }]
            }), 200

        except Exception as e:
            logger.error(f"Error fetching filing years: {str(e)}")
            return jsonify({"error": f"Failed to retrieve filing years: {str(e)}"}), 500
        


    def get_top_keywords(tfidf_matrix, feature_names, top_n=5):
        """
        Get top N keywords for each document in the TF-IDF matrix.
        """
        top_keywords = []
        for row in tfidf_matrix:
            top_indices = row.toarray().argsort()[0][-top_n:][::-1]
            top_words = [feature_names[i] for i in top_indices]
            top_keywords.append(top_words)
        return top_keywords

    @app.route('/api/topic_evolution', methods=['GET'])
    def topic_evolution():
        """
        Endpoint to analyze topic evolution in patent data.
        Stores data in patent_keywords only if keyword_df is not empty,
        analyzes topic evolution, and returns results as JSON.
        """
        try:
            # Query database for required columns
            patents = db.session.query(
                RawPatent.first_publication_number.label('first publication number'),
                RawPatent.title.label('Title'),
                RawPatent.first_filing_year.label('first filing year')
            ).all()
        
            if not patents:
                logger.info("No patents found in the database")
                return jsonify({"message": "No patents found in the database"}), 404
        
            # Create DataFrame from query results
            keyword_df = pd.DataFrame([row._asdict() for row in patents])
        
            # Drop rows where 'Title' is null
            keyword_df = keyword_df.dropna(subset=['Title'])
        
            if keyword_df.empty:
                logger.info("No patents with valid titles after filtering")
                return jsonify({"message": "No patents with valid titles"}), 404
        
            # Preprocess titles
            keyword_df['processed_title'] = keyword_df['Title'].apply(preprocess_text)
        
            # Compute TF-IDF
            vectorizer = TfidfVectorizer(max_df=0.85, min_df=2, ngram_range=(1, 3))
            tfidf_matrix = vectorizer.fit_transform(keyword_df['processed_title'])
            feature_names = vectorizer.get_feature_names_out()
        
            # Get top keywords for each patent
            top_keywords_list = get_top_keywords(tfidf_matrix, feature_names, top_n=5)
        
            # Store in patent_keywords only if keyword_df is not empty
            for i, row in keyword_df.iterrows():
                patent_keyword = PatentKeyword(
                    first_publication_number=row['first publication number'],
                    title=row['Title'],
                    first_filing_year=row['first filing year'],
                    keywords=top_keywords_list[i]
                )
                db.session.add(patent_keyword)
            db.session.commit()
        
            # Analyze topic evolution
            topic_evolution_data, windows_data , divergences, valid_years = analyze_topic_evolution(keyword_df)
        
            # Delete existing Window and Topic records
            db.session.query(Topic).delete()
            db.session.query(Window).delete()
            db.session.query(Divergence).delete()
            db.session.commit()
        
            # Insert new windows
            window_objects = []
            for win in windows_data:
                window = Window(start_year=win['start'], end_year=win['end'])
                db.session.add(window)
                window_objects.append(window)
            db.session.commit()
        
            # Map (start, end) to window objects
            window_dict = {(w.start_year, w.end_year): w for w in window_objects}
        
            # Insert topics
            for topic_data in topic_evolution_data:
                start = topic_data['start']
                end = topic_data['end']
                window = window_dict.get((start, end))
                if window:
                    topic_id_str = topic_data['topic_id']
                    topic_number = int(topic_id_str.split('-')[-1])
                    topic = Topic(
                        window_id=window.id,
                        topic_number=topic_number,
                        words=topic_data['words'],
                        weights=topic_data['weights']
                    )
                    db.session.add(topic)
            db.session.commit()
            
                   # Insert divergences
            for i in range(len(divergences)):
                from_year = int(valid_years[i])
                to_year = int(valid_years[i + 1])
                divergence_value = float(divergences[i])
                divergence_record = Divergence(from_year=from_year, to_year=to_year, divergence=divergence_value)
                db.session.add(divergence_record)
            db.session.commit()
        
            # Query stored data
            windows = Window.query.order_by(Window.start_year).all()
            topics = Topic.query.join(Window).order_by(Window.start_year, Topic.topic_number).all()
            divergences_query = Divergence.query.order_by(Divergence.from_year).all()

            window_list = [
                {
                    'start': w.start_year,
                    'end': w.end_year,
                    'years': list(range(w.start_year, w.end_year + 1))
                } for w in windows
            ]
        
            topic_evolution_list = [
                {
                    'start': t.window.start_year,
                    'end': t.window.end_year,
                    'topic_id': f"{t.window.start_year}-{t.topic_number}",
                    'words': t.words,
                    'weights': t.weights
                } for t in topics
            ]
            divergences_list = [
            {
                'from_year': d.from_year,
                'to_year': d.to_year,
                'divergence': d.divergence
            } for d in divergences_query
        ]
        
            # Query patent_keywords for response
            patent_keywords = PatentKeyword.query.all()
            patent_keywords_list = [
                {
                    'first_publication_number': pk.first_publication_number,
                    'title': pk.title,
                    'first_filing_year': pk.first_filing_year,
                    'keywords': pk.keywords
                } for pk in patent_keywords
            ]
        
            logger.info("Topic evolution analysis and keyword extraction completed successfully")
            return jsonify({
                "topic_evolution": topic_evolution_list,
                "windows": window_list,
                "patent_keywords": patent_keywords_list,
                "divergences": divergences_list
            }), 200
    
        except Exception as e:
            db.session.rollback()
            logger.error(f"Error in topic_evolution endpoint: {str(e)}")
            return jsonify({"error": f"Failed to process topic evolution: {str(e)}"}), 500


    
    @app.route('/api/automatic_topic_shift', methods=['GET'])
    def automatic_topic_shift():
        try:
            # Query all divergences, ordered by from_year
            divergences = Divergence.query.order_by(Divergence.from_year).all()
            if not divergences:
                return jsonify({"message": "No divergence data available"}), 200

            # Calculate threshold as the 80th percentile of divergence values
            divergence_values = [d.divergence for d in divergences]
            threshold = np.percentile(divergence_values, 80) if divergence_values else 0

            # Query all windows, ordered by start_year
            windows = Window.query.order_by(Window.start_year).all()

            # Format divergence data
            divergence_data = [
                {
                    "from_year": d.from_year,
                    "to_year": d.to_year,
                    "divergence": d.divergence
                }   
                for d in divergences
            ]

            # Format windows data
            window_list = [
                {
                    "start": w.start_year,
                    "end": w.end_year
                }
                for w in windows
            ]

            # Construct and return the JSON response
            response = {
                "divergence_data": divergence_data,
                "threshold": threshold,
                "windows": window_list
            }
            return jsonify(response), 200

        except Exception as e:
            return jsonify({"error": str(e)}), 500
        
        
    # @app.route('/api/evolving_word_clouds', methods=['GET'])
    # def evolving_word_clouds():
    #     try:
    #         # Query all windows, ordered by start year
    #         windows = Window.query.order_by(Window.start_year).all()
    #         result = []

    #         for window in windows:
    #             # Query topics associated with this window
    #             topics = Topic.query.filter_by(window_id=window.id).all()
            
    #             # Collect all words from topics in this window
    #             all_words = [word for topic in topics for word in topic.words]
            
    #             if not all_words:
    #                 continue
            
    #             # Count frequency of each word (replicates WordCloud's default behavior)
    #             word_freq = Counter(all_words)
            
    #             # Format word frequencies as a list of dictionaries
    #             word_freq_list = [{'word': word, 'frequency': freq} for word, freq in word_freq.items()]
            
    #             # Append window data to result
    #             result.append({
    #                 'start': window.start_year,
    #                 'end': window.end_year,
    #                 'word_frequencies': word_freq_list
    #             })
        
    #         return jsonify(result), 200
    
    #     except Exception as e:
    #         return jsonify({"error": str(e)}), 500

    @app.route('/api/weighted_word_clouds', methods=['GET'])
    def weighted_word_clouds():
        """
        Endpoint to generate word cloud data for each time window.
        Aggregates word weights across all topics within each window.
        Returns a JSON response with start/end years and word-weight pairs.
        """
        try:
            # Set up logging
            logger = logging.getLogger(__name__)

            # Query all windows, ordered by start year
            windows = Window.query.order_by(Window.start_year).all()
            result = []

            for window in windows:
                # Use defaultdict to accumulate weights for each word
                word_weight_dict = defaultdict(float)

                # Query all topics associated with this window
                topics = Topic.query.filter_by(window_id=window.id).all()

                for topic in topics:
                    # Validate that words and weights arrays have the same length
                    if len(topic.words) != len(topic.weights):
                        logger.warning(
                            f"Skipping topic {topic.id}: words and weights lengths mismatch "
                            f"({len(topic.words)} words, {len(topic.weights)} weights)"
                        )
                        continue

                    # Aggregate weights by summing for each word
                    for word, weight in zip(topic.words, topic.weights):
                        processed_word = preprocess_text(word.lower())
                        if processed_word and processed_word.strip():
                            word_weight_dict[word] += weight

                # Skip if no words were aggregated
                if not word_weight_dict:
                    logger.info(f"No valid words found for window {window.start_year}-{window.end_year}")
                    continue

                # Convert aggregated weights to a sorted list of dictionaries
                word_list = [
                    {"word": word, "weight": weight}
                    for word, weight in sorted(word_weight_dict.items(), key=lambda x: x[1], reverse=True)
                ]

                # Append window data to the result
                result.append({
                    "start": window.start_year,
                    "end": window.end_year,
                    "words": word_list
                })

            # Log successful completion
            logger.info(f"Generated word cloud data for {len(result)} windows")
            return jsonify(result), 200

        except Exception as e:
            logger.error(f"Error in weighted_word_clouds endpoint: {str(e)}")
            return jsonify({"error": str(e)}), 500
        
        
    @app.route('/classify_patents', methods=['POST'])
    def classify_patents():
        try:
            # 1) If the IPCClassification table is empty, load from CSV
            if IPCClassification.query.count() == 0:
                csv_path = 'classification_df.csv'  

               
                df = pd.read_csv(
                    csv_path,
                    sep=';',
                    usecols=['CPC Symbol', 'Classification Title']
                )

                # Rename columns so they match the model field names
                df.columns = ['cpc_symbol', 'classification_title']

                # Use cpc_symbol as the DataFrame index for easy lookup
                df.set_index('cpc_symbol', inplace=True)

                # Insert into the database
                for code, row in df.iterrows():
                    new_class = IPCClassification(
                        cpc_symbol=code,
                        classification_title=row['classification_title']
                    )
                    db.session.add(new_class)
                db.session.commit()
                
            if IPCFieldOfStudy.query.count() == 0:
                csv_path = 'IPC_to_fieldOfStudy.csv'  

                df = pd.read_csv(
                    csv_path,
                    sep=',',
                    usecols=lambda x: x in ['IPC', 'description', 'fields'] 
                )

                # Rename columns so they match the model field names
                df.columns = ['ipc','description', 'fields']

                # Use cpc_symbol as the DataFrame index for easy lookup
                df.set_index('ipc', inplace=True)

                # Insert into the database
                for code, row in df.iterrows():
                    new_field = IPCFieldOfStudy(
                        ipc=code,
                        fields=row['fields'],
                        description=row['description']
                    )
                    db.session.add(new_field)
                db.session.commit()

            # 2) Build a lookup DataFrame from whatever is currently in the table
            all_rows = IPCClassification.query.all()
            classification_dict = {
                row.cpc_symbol: row.classification_title
                for row in all_rows
            }
            classification_df = pd.DataFrame.from_dict(
                classification_dict,
                orient='index',
                columns=['classification_title']
            )
            classification_df.index.name = 'cpc_symbol'
            
            all_rows2 = IPCFieldOfStudy.query.all()
            field_of_study_dict = {
                row.ipc: {
                    'ipc': row.ipc,
                    'fields': row.fields,
                    'description': row.description
                }
                for row in all_rows2
            }

            field_of_study_df = pd.DataFrame.from_dict(
                field_of_study_dict,
                orient='index',
                columns=['ipc','fields', 'description']
            )
            field_of_study_df.index.name = 'ipc'
            field_of_study_df['fields'] = field_of_study_df['fields'].apply(lambda x: ast.literal_eval(x) if isinstance(x, str) else x)
            field_of_study_df['fields'] = field_of_study_df['fields'].apply(normalize_engineering)
            field_of_study_df['ipc'] = field_of_study_df['ipc'].str[:3].str.strip().str.upper()
            #ipc_mapping = field_of_study_df.set_index('IPC')['fields'].to_dict()
            # 3) Fetch unclassified patents and assign the “meaning” based on IPC codes
            patents = RawPatent.query.all()
            if not patents:
                return jsonify({"message": "No patents found to classify"}), 200

            updated_count = 0
            for patent in patents:
                if not patent.ipc:
                    continue
                # Cleanup step #1: remove leading/trailing braces & quotes:

                cleaned = patent.ipc.strip('"{}')
                # Cleanup step #2: split on commas, drop any empty strings
                
                ipc_list = [code.strip() for code in cleaned.split(',') if code.strip()]
                if not ipc_list:
                    continue
                
                meaning = get_ipc_meaning(ipc_list, classification_df)
                fields = get_patent_fields(ipc_list, field_of_study_df)  
                if isinstance(fields, str):
                    try:
                        fields = ast.literal_eval(fields)
                    except Exception:
                        fields = [fields]
                if meaning:
                # Check if this patent already exists in classified_patents
                    existing = ClassifiedPatent.query.get(patent.id)
                    if existing:
                        # Update existing record
                        existing.ipc_meaning = meaning
                        existing.fields = fields 
                    else:
                        # Create new classified patent record
                        classified_patent = ClassifiedPatent(
                            id=patent.id,
                            title=patent.title,
                            inventors=patent.inventors,
                            applicants=patent.applicants,
                            publication_number=patent.publication_number,
                            earliest_priority=patent.earliest_priority,
                            earliest_publication=patent.earliest_publication,
                            ipc=patent.ipc,
                            cpc=patent.cpc,
                            publication_date=patent.publication_date,
                            first_publication_date=patent.first_publication_date,
                            second_publication_date=patent.second_publication_date,
                            first_filing_year=patent.first_filing_year,
                            earliest_priority_year=patent.earliest_priority_year,
                            applicant_country=patent.applicant_country,
                            family_number=patent.family_number,
                            family_jurisdictions=patent.family_jurisdictions,
                            family_members=patent.family_members,
                            first_publication_number=patent.first_publication_number,
                            second_publication_number=patent.second_publication_number,
                            first_publication_country=patent.first_publication_country,
                            second_publication_country=patent.second_publication_country,
                            ipc_meaning=meaning,
                            fields=fields
                        )
                        db.session.add(classified_patent)
                    updated_count += 1

            db.session.commit()
            return jsonify({
                "message": f"Successfully processed {updated_count} patents",
                "classified_count": updated_count
            }), 200
        except Exception as e:
                db.session.rollback()
                return jsonify({"error": f"Failed to classify patents: {str(e)}"}), 500
        
    @app.route('/api/patent_field_trends', methods=['GET'])
    def patent_field_trends():
        """
        Returns JSON for plotting field trends over time.
        Query params:
        - top_n: number of top fields to return (default 5)
        - smoothing: rolling window size (default 3)
        """
        try:
            top_n = int(request.args.get('top_n', 5))
            smoothing = int(request.args.get('smoothing', 3))

            # Query classified patents with year and fields
            patents = ClassifiedPatent.query.with_entities(
                ClassifiedPatent.first_filing_year,
                ClassifiedPatent.fields
            ).filter(ClassifiedPatent.first_filing_year.isnot(None)).all()

            # Build DataFrame
            rows = []
            for year, fields in patents:
                # Ensure fields is a list
                if isinstance(fields, str):
                    try:
                        fields = ast.literal_eval(fields)
                    except Exception:
                        fields = [fields]
                for field in fields:
                    rows.append({'year': int(year), 'field': field})

            if not rows:
                return jsonify({"labels": [], "datasets": []})

            df = pd.DataFrame(rows)
            # Count papers per year/field
            field_counts = (
                df.groupby(['year', 'field'])
                .size()
                .reset_index(name='count')
            )
            # Pivot for trend lines
            pivot = field_counts.pivot(index='year', columns='field', values='count').fillna(0)
            # Smooth with rolling average
            pivot_smooth = pivot.rolling(window=smoothing, min_periods=1).mean()

            # Pick top N fields by total count
            top_fields = pivot.sum().sort_values(ascending=False).head(top_n).index
            labels = list(map(int, pivot_smooth.index))

            datasets = []
            for field in top_fields:
                datasets.append({
                    "label": field,
                    "data": [int(v) for v in pivot_smooth[field].values]
                })

            return jsonify({
                "labels": labels,
                "datasets": datasets
            }), 200

        except Exception as e:
            return jsonify({"error": str(e)}), 500
        

    @app.route('/api/cooccurrence_trends', methods=['GET'])
    def cooccurrence_trends():
        """
        Returns JSON data for plotting co-occurrence trends of technology keyword pairs.
        Removes English and French stopwords from titles and ensures term pairs are not identical.
        """
        try:
            # Fetch patent data from the database
            query = 'SELECT "Publication number", "Title", "first_filing_year" FROM raw_patents'
            df = pd.read_sql(query, engine)
            df = df.dropna(subset=['Title', 'first_filing_year'])

            # Remove stopwords from Title
            df['Title'] = df['Title'].apply(clean_text_remove_stopwords)

            # Remove empty titles after cleaning
            df = df[df['Title'].str.strip().astype(bool)]

            # Prepare the grouped DataFrame as required by track_cooccurrence_trends
            grouped = df.groupby("first_filing_year")["Title"].apply(lambda texts: " ".join(texts)).reset_index()

            # Run the co-occurrence trend analysis
            cooc_trends = track_cooccurrence_trends(
                grouped,
                time_col='first_filing_year',
                text_col='Title',
                window_size=5,
                min_count=10
            )

            # Filter out pairs where both terms are the same (case-insensitive)
            cooc_trends = cooc_trends[cooc_trends['term1'].str.lower() != cooc_trends['term2'].str.lower()]

            # Get top emerging and declining combinations
            emerging_tech = cooc_trends[
                (cooc_trends.slope > 0) & (cooc_trends.p_value < 0.05)
            ].sort_values('slope', ascending=False).head(5)

            declining_tech = cooc_trends[
                (cooc_trends.slope < 0) & (cooc_trends.p_value < 0.05)
            ].sort_values('slope').head(5)

            # Prepare data for plotting
            def prepare_plot_data(df):
                plot_data = []
                for _, row in df.iterrows():
                    for year, freq in row['frequency_history']:
                        plot_data.append({
                            'year': int(year),
                            'frequency': int(freq),
                            'term_pair': f"{row['term1']} & {row['term2']}"
                        })
                return plot_data

            response = {
                "emerging": prepare_plot_data(emerging_tech),
                "declining": prepare_plot_data(declining_tech),
                "emerging_pairs": [
                    {
                        "term1": row['term1'],
                        "term2": row['term2'],
                        "slope": row['slope'],
                        "total_count": row['total_count']
                    }
                    for _, row in emerging_tech.iterrows()
                ],
                "declining_pairs": [
                    {
                        "term1": row['term1'],
                        "term2": row['term2'],
                        "slope": row['slope'],
                        "total_count": row['total_count']
                    }
                    for _, row in declining_tech.iterrows()
                ]
            }
            return jsonify(response), 200

        except Exception as e:
            return jsonify({"error": str(e)}), 500
        
        
        
    @app.route('/api/analyze_applicants', methods=['POST'])
    def analyze_applicants():
        # Load patent data from DB
        patents = RawPatent.query.all()
        data = []
        for patent in patents:
            data.append({
                'first applicant': patent.applicants.split('\n')[0] if patent.applicants else None,
                'second applicant': patent.applicants.split('\n')[1] if patent.applicants and '\n' in patent.applicants else None,
                'Inventors': patent.inventors
            })
        df = pd.DataFrame(data)
        # Fill missing columns if needed
        if 'Inventors' not in df.columns:
            df['Inventors'] = None

        # Generate DataFrames
        applicants_df = get_applicants_df(df)
        applicant_type_df = get_applicant_type_df(applicants_df)

        # Store applicants_df
        Applicant.query.delete()
        for _, row in applicants_df.iterrows():
            db.session.add(Applicant(applicant_name=row['Applicants'], applicant_type=row['Applicant Type']))
        db.session.commit()

        # Store applicant_type_df
        ApplicantType.query.delete()
        for _, row in applicant_type_df.iterrows():
            db.session.add(ApplicantType(applicant_type=row['Applicant Type'], percentage=row['Percentage']))
        db.session.commit()

        return jsonify({"status": "success"})
    
    
    @app.route('/api/applicant_type_summary', methods=['GET'])
    def applicant_type_summary():
        """
        Returns applicant type distribution for visualization.
    Example response:
    {
        "labels": ["Company - Incorporated/Corporation", "University/Research Institution", ...],
        "percentages": [45.2, 30.1, ...],
        "details": [
            {"applicant_name": "Toyota Motor Corp", "applicant_type": "Company - Incorporated/Corporation"},
            ...
        ]
    }
     """
        # Query summary table
        summary = ApplicantType.query.order_by(ApplicantType.percentage.desc()).all()
        labels = [row.applicant_type for row in summary]
        percentages = [row.percentage for row in summary]

        # Optionally, provide a detailed list (for tooltips, drill-down, etc.)
        applicants = Applicant.query.all()
        details = [
            {"applicant_name": a.applicant_name, "applicant_type": a.applicant_type}
            for a in applicants
        ]

        return jsonify({
            "labels": labels,
            "percentages": percentages,
            "details": details
        })

    def load_applicants_df() -> pd.DataFrame:
        """
        Reads the entire `applicants` table into a DataFrame.
        """
        # this issues: SELECT * FROM applicants
        df = pd.read_sql_table(
            table_name='applicants',
            con=db.engine,
            columns=['id', 'applicant_name', 'applicant_type']
        )
        return df


    @app.route('/api/innovation_cycle', methods=['GET'])
    def innovation_cycle():
        # 1) Pull all “Applicants” strings from the DB
        df = load_applicants_df()
        df = df.rename(columns={'applicant_name': 'Applicants'})
        cycle = get_innovation_cycle(df, top_n=10)
        # 5) Return as JSON
        return jsonify(f"{cycle:.2f}"), 200
    
    
    @app.route('/api/coapplicant_rate', methods=['GET'])
    def coapplicant_rate():
        try:
        
            # Load applicant data
            query = """
            SELECT "Applicants" 
            FROM raw_patents 
            WHERE "Applicants" IS NOT NULL
            """
            df = pd.read_sql(query, engine)
            if df.empty:
                return jsonify({
                    "coapplicant_rate": 0.0,
                    "total_applications": 0,
                    "coapplicant_count": 0
                })
        
            # Split applicants into first and second
            split_applicants = df['Applicants'].str.split('\n', n=1, expand=True)
            df['first_applicant'] = split_applicants[0]
            df['second_applicant'] = split_applicants[1]
        
            # Calculate co-applicant rate
            total_applications = len(df)
            coapplicant_count = df['second_applicant'].notnull().sum()
            coapplicant_rate = (coapplicant_count / total_applications) * 100
        
            return jsonify({
                "coapplicant_rate": float(round(coapplicant_rate, 2)),
                "total_applications": int(total_applications),
                "coapplicant_count": int(coapplicant_count)
            })
        
        except Exception as e:
            app.logger.error(f"Error in coapplicant_rate: {str(e)}")
            return jsonify({"error": str(e)}), 500
        
        
        
    @app.route('/api/growth_rate', methods=['GET'])  
    def growth_rate():
        df = pd.read_sql_table('raw_patents', con=engine, columns=['first_filing_year'])
        growth_rate = patent_growth_summary(df)
        return jsonify({"growth_rate": growth_rate}), 200

                   
            
    @app.route('/api/market_cost', methods=['POST'])
    def update_costs():
        try:
            # Step 1: Ensure Cost table is populated
            if Cost.query.count() == 0:
                csv_path = os.path.join(os.path.dirname(__file__), 'corrected-patent-cost-data.csv')
                if not os.path.exists(csv_path):
                    return jsonify({"error": "CSV file not found"}), 404
                cost = pd.read_csv(csv_path)
                cost.columns = [
                'Country',   
                'Years 0.0-1.5',
                'Years 2.0-4.5',
                'Years 5.0-9.5',
                'Years 10.0-14.5',
                'Years 15.0-20.0',
                'Total Cost (US$)'
                ]
                updated_cost = update_cost_df(cost)
                rename_dict = {
                'Years 0.0-1.5': 'Years_0_1_5',
                'Years 2.0-4.5': 'Years_2_4_5',
                'Years 5.0-9.5': 'Years_5_9_5',
                'Years 10.0-14.5': 'Years_10_14_5',
                'Years 15.0-20.0': 'Years_15_20',
                'Total Cost (US$)': 'Total_cost'
                }
                updated_cost.rename(columns=rename_dict, inplace=True)
                for _, row in updated_cost.iterrows():
                    new_cost = Cost(
                        Country=row['Country'],
                        Years_0_1_5=row['Years_0_1_5'],
                        Years_2_4_5=row['Years_2_4_5'],
                        Years_5_9_5=row['Years_5_9_5'],
                        Years_10_14_5=row['Years_10_14_5'],
                        Years_15_20=row['Years_15_20'],
                        Total_cost=row['Total_cost']
                    )
                    db.session.add(new_cost)
                db.session.commit()

            # Step 2: Load cost data
            cost_df = pd.read_sql('SELECT * FROM costs', db.engine)

            # Step 3: Fetch patent data into a DataFrame
            query = 'SELECT "earliest_priority_year", "first_publication_country" FROM raw_patents'
            patent_df = pd.read_sql(query, db.engine)

            # Step 4: Calculate patent age
            patent_df = calculate_age(patent_df)

            # Step 5: Assign cost to each patent
            patent_df['cost'] = patent_df.apply(lambda row: assign_cost(row, cost_df), axis=1)

            # Step 6: Clear existing patent_costs and save new data
            db.session.query(PatentCost).delete()
            db.session.commit()
            for _, row in patent_df.iterrows():
                if pd.notnull(row['Patent Age']) and pd.notnull(row['cost']):
                    patent_cost = PatentCost(
                        patent_age=str(row['Patent Age']),
                        country=row['first_publication_country'],
                        cost=row['cost']
                    )
                    db.session.add(patent_cost)
            db.session.commit()

            return jsonify({"message": "Patent costs updated successfully"}), 200

        except Exception as e:
            db.session.rollback()
            app.logger.error(f"Error in update_costs: {str(e)}")
            return jsonify({"error": str(e)}), 500
        
        
        






    @app.route('/api/originality_rate/fetch', methods=['POST'])
    def fetch_originality_data():
        request_id = str(uuid.uuid4())
        print(f"[{datetime.utcnow()}] Request ID: {request_id} - Fetch started")

        try:
            limit = request.args.get('limit', default=50, type=int)
            with app.app_context():
                db.create_all()
            patents_added = fetch_and_store_originality_data(limit=limit)
            # How many patents are in the DB now (valid for originality-rate calc)?
            _, _, valid_patents = calculate_originality_rate_from_db()

            print(f"[{datetime.utcnow()}] Request ID: {request_id} - Fetch completed")
            return jsonify({
                "status": "success",
                "message": f"Originality data fetched and stored (limit={limit}).",
                "patents_added": patents_added,
                "processed_patents": valid_patents, 
                "request_id": request_id
            }), 200
        except Exception as e:
            print(f"[{datetime.utcnow()}] Request ID: {request_id} - Error: {str(e)}")
            return jsonify({"error": str(e), "request_id": request_id}), 500

    @app.route('/api/originality_rate', methods=['GET'])
    def originality_rate():
        try:
            originality_rate, total_patents, n_valid = calculate_originality_rate_from_db()
            if originality_rate is None:
                return jsonify({
                    "error": f"Insufficient data: only {n_valid} valid patents. Need at least 30."
                }), 400
            return jsonify({
                "originality_rate": originality_rate,
                "total_patents": total_patents,
                "valid_patents": n_valid
            }), 200
        except Exception as e:
            app.logger.error(f"Error in originality_rate: {str(e)}")
            return jsonify({"error": str(e)}), 500
        
        
        
        
    @app.route('/api/market_metrics', methods=['GET'])
    def market_metrics():
        try:
            # Load all patents and cost data from the database
            patents = [p.to_dict() for p in RawPatent.query.all()]
            costs = [c.to_dict() for c in Cost.query.all()]
            import pandas as pd
            patents_df = pd.DataFrame(patents)
            cost_df = pd.DataFrame(costs)
            # Defensive: If either is empty, return zeros
            if patents_df.empty or cost_df.empty:
                return jsonify({
                'market_value': 0.0,
                'market_rate': 0.0,
                'mean_value': 0.0
            }), 200
            # Ensure 'Patent Age' column exists
            if 'Patent Age' not in patents_df.columns:
                from datetime import datetime
                current_year = datetime.now().year
                patents_df['Patent Age'] = current_year - patents_df['earliest_priority_year']
            # Compute metrics
            market_value, market_rate, mean_value = get_market_metrics(patents_df, cost_df)
            return jsonify({
                'market_value': float(market_value),
                'market_rate': float(market_rate),
                'mean_value': float(mean_value)
            }), 200
        except Exception as e:
            return jsonify({'error': str(e)}), 500



            
            
             





#@app.route('/api/family_members/API', methods=['POST']) 

    
    # @app.route('/api/family', methods=['POST'])
    # def populate_family_members():
    # # Get JSON data from the request
    #     CONSUMER_KEY = os.getenv("CONSUMER_KEY_3").strip()
    #     CONSUMER_SECRET = os.getenv("CONSUMER_SECRET_3").strip()
    #     CONSUMER_KEY1 = os.getenv("CONSUMER_KEY_2").strip()
    #     CONSUMER_SECRET1 = os.getenv("CONSUMER_SECRET_2").strip()

    #     try:
    #     # Fetch data from database
    #         query = 'SELECT * FROM raw_patents'
    #         df = pd.read_sql(query, engine)

    #     # Data cleaning
    #         df.rename(columns={
    #             'Titre': 'Title',
    #             'Inventeurs': 'Inventors',
    #             'Demandeurs': 'Applicants',
    #             'Numéro de publication': 'Publication number',
    #             'Priorité la plus ancienne': 'Earliest priority',
    #             'CIB': 'IPC',
    #             'CPC': 'CPC',
    #             'Date de publication': 'Publication date',
    #             'Publication la plus ancienne': 'Earliest publication',
    #             'Numéro de famille': 'Family number'
    #         }, inplace=True)

    #         # Split 'Publication date' into two columns using regex
    #         df['first publication date'] = df['Publication date'].str.extract(r'^(\S+)', expand=False)
    #         df['second publication date'] = df['Publication date'].str.extract(r'^\S+\s+(.*)', expand=False)


    #         df['second publication date'] = df['second publication date'].str.strip('\r\n')
        
    #         df[['first publication number', 'second publication number']] = df['Publication number'].str.split(' ', n=1, expand=True)
    #         df['second publication number'] = df['second publication number'].str.strip('\r\n')
        
    #         if 'Unnamed: 11' in df.columns:
    #             df.drop(columns=['Unnamed: 11', 'Publication date'], inplace=True)
        
    #         df['family number'] = pd.to_numeric(df['Family number'], errors='coerce')
    #         df.rename(columns={'Family number': 'family number'}, inplace=True)
            
    #         # Calculate the number of rows for each part
    #         n = len(df) // 3

    #     # Split the DataFrame into three parts
    #         df1 = df.iloc[:n].copy()       # First part
    #         df2 = df.iloc[n:2*n].copy()    # Second part
    #         df3 = df.iloc[2*n:].copy()     # Third part

    #     # Process df1
    #         df1 = process_dataframe_parallel(df1, 'first publication number', max_workers=4)
    #         print('num of null values df1 :', df1['family_members'].isnull().sum(), 'number of empty arrays : ', df1['family_jurisdictions'].apply(lambda x: isinstance(x, list) and len(x) == 0).sum())
    #     #process df2
    #         search_keywords = fetch_last_search_keywords_from_db()
    #         if 'family_members' not in df2.columns: 
    #             df['family_members'] = None
    #         #split the dataframe into 3 parts 
    #         indices = df2.index.tolist()
    #         n = len(indices)
    #         part_size = n // 3
    #         remainder = n % 3
    #         parts= []
    #         start = 0
    #         for i in range(3) : 
    #             if i < remainder:
    #                 end = start + part_size + 1
    #             else:
    #                 end = start + part_size 
    #             parts.append(indices[start:end])
    #             start = end 
    #         #create three threads , each with ist own patentsSearch instance 
    #         threads = []
    #         for part in parts : 
    #             thread = threading.Thread(target=process_rows, args=(df2, part,search_keywords,False))
    #             threads.append(thread)
    #         for thread in threads:
    #             thread.start()
    #         for thread in threads:
    #             thread.join()
    #         print('num of null values df2 :', df2['family_members'].isnull().sum(), 'number of empty arrays : ', df2['family_jurisdictions'].apply(lambda x: isinstance(x, list) and len(x) == 0).sum())
    #         #add the processing of the rows of family members here

    #     #process df3
    #         CONSUMER_KEY = os.getenv("CONSUMER_KEY").strip()
    #         CONSUMER_SECRET = os.getenv("CONSUMER_SECRET").strip()
    #         CONSUMER_KEY1 = os.getenv("CONSUMER_KEY_1").strip()
    #         CONSUMER_SECRET1 = os.getenv("CONSUMER_SECRET_1").strip()
            
    #         df3 = process_dataframe_parallel(df3, 'first publication number', max_workers=4)
    #         print('num of null values df3 :', df3['family_members'].isnull().sum(), 'number of empty arrays : ', df3['family_jurisdictions'].apply(lambda x: isinstance(x, list) and len(x) == 0).sum())
        
    #     #merging the dataframes
    #         df = pd.concat([df1, df2, df3], ignore_index=True)
    #         df['family_members'] = df['family_members'].apply(lambda x: x if isinstance(x, list) else [])
    #         df['family_jurisdictions'] = df['family_jurisdictions'].apply(lambda x: x if isinstance(x, list) else [])
    #         # Ensure the columns exist in the DataFrame
    #         ensure_columns_exist(df, ['family_members', 'family_jurisdictions'])
    #         print('num of null values df :', df['family_members'].isnull().sum(), 'number of empty arrays : ', df['family_jurisdictions'].apply(lambda x: isinstance(x, list) and len(x) == 0).sum())
            
    #     # prepare updates for the database 
    #         updates = [
    #             {'id':row['id'],'first publication number' : row['first publication number'], 'jurisdictions' : json.dumps(row['family_jurisdictions']), 'members' : json.dumps(row['family_members'])
    #              } for index, row in df.iterrows()
    #         ]
    #         #update the database with the new columns 
    #         try:
    #             engine.execute(
    #                 text("UPDATE raw_patents SET family_jurisdictions = :jurisdictions, family_members = :members WHERE id = :id"),
    #                 updates
    #             )
    #         except Exception as e:
    #             return jsonify({"error": f"failed to update database: {str(e)}"}), 500


    #     # Prepare response
    #         results = df[['first publication number', 'family_jurisdictions', 'family_members']].to_dict(orient='records')
    #         empty_arrays_count = df['family_jurisdictions'].apply(lambda x: isinstance(x, list) and len(x) == 0).sum()
    #         null_count = df['family_jurisdictions'].isnull().sum()

    #         return jsonify({
    #             "results": results,
    #             "statistics": {
    #             "empty_jurisdictions_count": int(empty_arrays_count),
    #             "null_jurisdictions_count": int(null_count),
    #             "total_processed": len(df)
    #         }
    #     })

    #     except sqlalchemy.exc.SQLAlchemyError as e:
    #         return jsonify({"error": f"Database error: {str(e)}"}), 500
    #     except Exception as e:
    #         return jsonify({"error": f"Processing failed: {str(e)}"}), 500
    
    
    
    @app.route('/api/family', methods=['POST'])
    def populate_family_members():
        try:
            # Load API credentials from environment variables
            CONSUMER_KEY = os.getenv("CONSUMER_KEY")
            CONSUMER_SECRET = os.getenv("CONSUMER_SECRET")
            CONSUMER_KEY1 = os.getenv("CONSUMER_KEY1")
            CONSUMER_SECRET1 = os.getenv("CONSUMER_SECRET1")

            # Fetch data from database
            query = 'SELECT * FROM raw_patents'
            df = pd.read_sql(query, engine)

            # Data cleaning
            df.rename(columns={
                'Titre': 'Title',
                'Inventeurs': 'Inventors',
                'Demandeurs': 'Applicants',
                'Numéro de publication': 'Publication number',
                'Priorité la plus ancienne': 'Earliest priority',
                'CIB': 'IPC',
                'CPC': 'CPC',
                'Date de publication': 'Publication date',
                'Publication la plus ancienne': 'Earliest publication',
                'Numéro de famille': 'Family number'
            }, inplace=True)

            df['first publication date'] = df['Publication date'].str.extract(r'^(\S+)', expand=False)
            df['second publication date'] = df['Publication date'].str.extract(r'^\S+\s+(.*)', expand=False)
            df['second publication date'] = df['second publication date'].str.strip('\r\n')
            df[['first publication number', 'second publication number']] = df['Publication number'].str.split(' ', n=1, expand=True)
            df['second publication number'] = df['second publication number'].str.strip('\r\n')
            if 'Unnamed: 11' in df.columns:
                df.drop(columns=['Unnamed: 11', 'Publication date'], inplace=True)
            df['family number'] = pd.to_numeric(df['Family number'], errors='coerce')
            df.rename(columns={'Family number': 'family number'}, inplace=True)

            # Split the DataFrame into three parts
            n = len(df) // 3
            df1 = df.iloc[:n].copy()
            df2 = df.iloc[n:2*n].copy()
            df3 = df.iloc[2*n:].copy()

            # Process df1 using API with parallel processing
            df1 = process_dataframe_parallel(df1, 'first publication number', CONSUMER_KEY, CONSUMER_SECRET, max_workers=4)

            # Process df2 using web scraping with threading
            search_keywords = fetch_last_search_keywords_from_db()
            if 'family_members' not in df2.columns:
                df2['family_members'] = None
            indices = df2.index.tolist()
            part_size = len(indices) // 3
            parts = [indices[i:i + part_size] for i in range(0, len(indices), part_size)]
            threads = []
            for part in parts:
                thread = threading.Thread(target=process_rows, args=(df2, part, search_keywords, True))
                threads.append(thread)
                time.sleep(random.uniform(0.5, 1.5))  # Stagger thread starts
            for thread in threads:
                thread.start()
            for thread in threads:
                thread.join()

            # Process df3 using API with parallel processing and different credentials
            df3 = process_dataframe_parallel(df3, 'first publication number', CONSUMER_KEY1, CONSUMER_SECRET1, max_workers=4)

            # Merge the DataFrames
            df = pd.concat([df1, df2, df3], ignore_index=True)
            df['family_members'] = df['family_members'].apply(lambda x: x if isinstance(x, list) else [])
            df['family_jurisdictions'] = df['family_jurisdictions'].apply(lambda x: x if isinstance(x, list) else [])

            # Prepare updates for the database
            updates = [
                {'id': row['id'], 'jurisdictions': json.dumps(row['family_jurisdictions']), 'members': json.dumps(row['family_members'])}
                for _, row in df.iterrows()
            ]

            # Update the database
            with engine.connect() as connection:
                connection.execute(
                    text("UPDATE raw_patents SET family_jurisdictions = :jurisdictions, family_members = :members WHERE id = :id"),
                updates
            )
            connection.commit()

            # Prepare response
            results = df[['first publication number', 'family_jurisdictions', 'family_members']].to_dict(orient='records')
            empty_arrays_count = df['family_jurisdictions'].apply(lambda x: isinstance(x, list) and len(x) == 0).sum()
            null_count = df['family_jurisdictions'].isnull().sum()

            return jsonify({
                "results": results,
                "statistics": {
                    "empty_jurisdictions_count": int(empty_arrays_count),
                    "null_jurisdictions_count": int(null_count),
                    "total_processed": len(df)
                }
            })
        except Exception as e:
            return jsonify({"error": f"Processing failed: {str(e)}"}), 500

    API_CREDENTIALS = load_api_credentials()
    TOKEN_URL = "https://ops.epo.org/3.2/auth/accesstoken"
    BASE_URL = "https://ops.epo.org/3.2/rest-services"
    # One token per credential
    TOKEN_CACHE = [{'token': None, 'expiry': 0} for _ in API_CREDENTIALS]
    TOKEN_LOCKS = [threading.Lock() for _ in API_CREDENTIALS]
    

    @app.route('/api/family_members/ops', methods=['POST'])
    def update_family_members_ops():
        """
        Update family_members and family_jurisdictions in raw_patents using EPO OPS API,
        with batch commits, key rotation, and optional scraping fallback.
        """
        # Query all patents (adjust query as needed, e.g. filter incomplete only)
        patents = RawPatent.query.all()
        batch_size = 100
        updates, failed = [], []
        progress = tqdm(total=len(patents), desc="Updating Patents", unit="patent")
        db_session = scoped_session(db.session)

        for idx, patent in enumerate(patents):
            pub_number = getattr(patent, 'first_publication_number', None)
            fam_number = getattr(patent, 'family_number', None)
            if not validate_patent_number(pub_number):
                failed.append({'id': patent.id, 'reason': 'invalid publication number'})
                progress.update(1)
                continue

            cred_idx = idx % len(API_CREDENTIALS)  # rotate API keys
            fam, err = fetch_family_data_api(pub_number, cred_idx)

            # Optional: Scraping fallback
            if fam is None:
                fam, err2 = fetch_family_data_scrape(pub_number, fam_number)
                if fam is None:
                    failed.append({'id': patent.id, 'reason': f"{err} | {err2}"})
                    progress.update(1)
                    continue

            # Update model (assuming JSON fields)
            patent.family_jurisdictions = fam['jurisdictions']
            patent.family_members = fam['family_members']
            updates.append({'id': patent.id, 'jurisdictions': fam['jurisdictions'], 'members': fam['family_members']})

            # Batch commit
            if len(updates) % batch_size == 0:
                try:
                    db_session.commit()
                    updates.clear()  # Only clear on successful commit
                except Exception as e:
                    db_session.rollback()
                    failed.append({'batch_failed_at': idx, 'error': str(e)})
            progress.update(1)
            time.sleep(1.2)  # Respect EPO rate limits

        # Final commit for any remaining updates
        try:
            db_session.commit()
        except Exception as e:
            db_session.rollback()
            return jsonify({'success': False, 'error': f'Database commit failed: {str(e)}'}), 500
        finally:
            db_session.remove()

        progress.close()
        return jsonify({
            'success': True,
            'updated': len(patents) - len(failed),
            'failed': failed[:100],  # truncate in case of huge jobs
            'sample_update': updates[0] if updates else None
        })
    
    
    
    # @app.route('/api/family_members/ops', methods=['POST'])
    # def update_family_members_ops():
    #     """
    #     Update family_members and family_jurisdictions in raw_patents using EPO OPS API.
    #     """
    #     import requests
    #     import time
    #     from urllib.parse import quote
    #     import os
    #     from dotenv import load_dotenv
    #     load_dotenv()
    #     CONSUMER_KEY = os.getenv("CONSUMER_KEY").strip()
    #     CONSUMER_SECRET = os.getenv("CONSUMER_SECRET").strip()
    #     TOKEN_URL = "https://ops.epo.org/3.2/auth/accesstoken"
    #     BASE_URL = "https://ops.epo.org/3.2/rest-services"
    #     TOKEN = None
    #     TOKEN_EXPIRY = 0

    #     def get_access_token():
    #         nonlocal TOKEN, TOKEN_EXPIRY
    #         if TOKEN and time.time() < TOKEN_EXPIRY:
    #             return TOKEN
    #         data = {
    #             "grant_type": "client_credentials",
    #             "client_id": CONSUMER_KEY,
    #             "client_secret": CONSUMER_SECRET
    #         }
    #         headers = {"Content-Type": "application/x-www-form-urlencoded"}
    #         response = requests.post(TOKEN_URL, data=data, headers=headers, timeout=15)
    #         response.raise_for_status()
    #         TOKEN = response.json()["access_token"]
    #         TOKEN_EXPIRY = time.time() + 3500  # ~58 min
    #         return TOKEN

    #     def validate_patent_number(patent):
    #         if not patent or len(str(patent).strip()) < 4:
    #             return False
    #         return True

    #     def extract_jurisdictions_and_members(data):
    #         try:
    #             jurisdictions = set()
    #             family_members = []
    #             world_data = data.get('ops:world-patent-data', {})
    #             patent_family = world_data.get('ops:patent-family', {})
    #             members = patent_family.get('ops:family-member', [])
    #             if isinstance(members, dict):
    #                 members = [members]
    #             for member in members:
    #                 pub_ref = member.get('publication-reference', {})
    #                 docs = pub_ref.get('document-id', [])
    #                 if isinstance(docs, dict):
    #                     docs = [docs]
    #                 for doc in docs:
    #                     if doc.get('@document-id-type') == 'docdb':
    #                         country = doc.get('country')
    #                         if isinstance(country, dict):
    #                             country = country.get('$')
    #                         doc_number = doc.get('doc-number')
    #                         if isinstance(doc_number, dict):
    #                             doc_number = doc_number.get('$')
    #                         kind = doc.get('kind')
    #                         if isinstance(kind, dict):
    #                             kind = kind.get('$')
    #                         if country and doc_number and kind:
    #                             jurisdictions.add(country)
    #                             family_members.append(f"{country}{doc_number}{kind}")
    #             return {
    #                 'jurisdictions': sorted(jurisdictions),
    #                 'family_members': sorted(set(family_members))
    #             }
    #         except Exception as e:
    #             print(f"Error parsing response: {e}")
    #             return {'jurisdictions': None, 'family_members': None}

    #     # Query all patents
    #     patents = RawPatent.query.all()
    #     updates = []
    #     failed = []
    #     request_delay = 1.2
    #     for patent in patents:
    #         pub_number = patent.first_publication_number
    #         if not validate_patent_number(pub_number):
    #             failed.append({'id': patent.id, 'reason': 'invalid publication number'})
    #             continue
    #         try:
    #             token = get_access_token()
    #             url = f"{BASE_URL}/family/publication/docdb/{quote(str(pub_number))}"
    #             headers = {
    #                 "Authorization": f"Bearer {token}",
    #                 "Accept": "application/json"
    #             }
    #             response = requests.get(url, headers=headers, timeout=15)
    #             if response.status_code == 403:
    #                 failed.append({'id': patent.id, 'reason': '403 forbidden'})
    #                 continue
    #             if response.status_code == 404:
    #                 failed.append({'id': patent.id, 'reason': '404 not found'})
    #                 continue
    #             response.raise_for_status()
    #             data = response.json()
    #             fam = extract_jurisdictions_and_members(data)
    #             # Update ORM object
    #             patent.family_jurisdictions = fam['jurisdictions']
    #             patent.family_members = fam['family_members']
    #             updates.append({'id': patent.id, 'jurisdictions': fam['jurisdictions'], 'members': fam['family_members']})
    #         except Exception as e:
    #             failed.append({'id': patent.id, 'reason': str(e)})
    #         time.sleep(request_delay)
    #     try:
    #         db.session.commit()
    #     except Exception as e:
    #         db.session.rollback()
    #         return jsonify({'success': False, 'error': f'Database commit failed: {str(e)}'}), 500
    #     return jsonify({
    #         'success': True,
    #         'updated': len(updates),
    #         'failed': failed,
    #         'sample_update': updates[0] if updates else None
    #     })









    @app.route('/api/family_member_by_country', methods=['GET'])
    def family_member_by_country():
        try:
            # Query all family_members from RawPatent
            patents = RawPatent.query.with_entities(RawPatent.family_members).all()
            # Convert to DataFrame
            import pandas as pd
            family_df = pd.DataFrame([{'family_members': p[0]} for p in patents])
            # Defensive: drop missing or empty
            family_df = family_df.dropna(subset=['family_members'])
            family_df = family_df[family_df['family_members'].apply(lambda x: isinstance(x, list) and len(x) > 0)]
            # Extract country codes for each family
            from cleaners import extract_country_codes_from_list
            family_df['country_codes'] = family_df['family_members'].apply(extract_country_codes_from_list)
            # Use analysis function
            family_counts_df = get_family_counts_by_country(family_df)
            # Prepare JSON for plotting
            result = {
                'labels': family_counts_df['country_code'].tolist(),
                'datasets': [{
                    'label': 'Number of Family Members',
                    'data': family_counts_df['member_count'].tolist(),
                    'backgroundColor': 'rgba(93, 164, 214, 0.7)'
                }]
            }
            return jsonify(result), 200
        except Exception as e:
            return jsonify({'error': str(e)}), 500



    @app.route('/api/top_10_patent_applicants', methods=['GET'])
    def top_10_patent_applicants():
        try:
            # Load applicants table into DataFrame
            df = pd.read_sql_table(
                table_name='applicants',
                con=db.engine,
                columns=['id', 'applicant_name']
            )
            # Count top 10 applicants
            top10 = df['applicant_name'].value_counts().head(10)
            result = {
                'labels': top10.index.tolist(),
                'datasets': [{
                    'label': 'Patent Count',
                    'data': top10.values.tolist(),
                    'backgroundColor': 'rgba(255, 99, 132, 0.7)'
                }]
            }
            return jsonify(result), 200
        except Exception as e:
            return jsonify({'error': str(e)}), 500

    @app.route('/api/applicant_collaboration_network', methods=['GET'])
    def applicant_collaboration_network():
        try:
            # Load patent data (Applicants column) from RawPatent
            patents = RawPatent.query.with_entities(RawPatent.applicants).all()
            import pandas as pd
            patent_df = pd.DataFrame([{'Applicants': p[0]} for p in patents])
            # Load applicant data from applicants table
            applicant_df = pd.read_sql_table(
                table_name='applicants',
                con=db.engine,
                columns=['applicant_name', 'applicant_type']
            )
            
            network_data = extract_applicant_collaboration_network(patent_df, applicant_df)
            return jsonify(network_data), 200
        except Exception as e:
            return jsonify({'error': str(e)}), 500

    @app.route('/api/top_ipc_codes', methods=['GET'])
    def top_ipc_codes():
        try:
            # Load IPC data from RawPatent
            patents = RawPatent.query.with_entities(RawPatent.ipc).all()

            ipc_list = []
            for p in patents:
                if p[0]:
                    # IPC may be a string with multiple codes separated by ';' or ','
                    codes = str(p[0]).replace(',', ';').split(';')
                    ipc_list.extend([c.strip() for c in codes if c.strip()])
            # Extract main IPC (before '/')
            main_ipc = [code.split('/')[0] for code in ipc_list if '/' in code]
            main_ipc_counts = Counter(main_ipc).most_common(10)
            main_ipc_df = pd.DataFrame(main_ipc_counts, columns=['Main IPC', 'Count'])

            # --- IPC meaning integration ---
            
            # Use the actual top IPC codes (main_ipc_df['Main IPC'])
            ipc_meanings = get_ipc_codes_meanings(main_ipc_df['Main IPC'].tolist())
            # Build a mapping for the frontend: list of dicts with code, count, title, explanation
            ipc_info = []
            for _, row in main_ipc_df.iterrows():
                code = row['Main IPC']
                count = row['Count']
                meaning = ipc_meanings.get(code, {})
                ipc_info.append({
                    'ipc_code': code,
                    'count': count,
                    'symbol': meaning.get('symbol', ''),
                    'title': meaning.get('title', ''),
                    'explanation': meaning.get('explanation', '')
                })
            result = {
                'labels': main_ipc_df['Main IPC'].tolist(),
                'datasets': [{
                    'label': 'Patent Count',
                    'data': main_ipc_df['Count'].tolist(),
                    'backgroundColor': 'rgba(153, 50, 204, 0.7)'
                }],
                'ipc_info': ipc_info
            }
            return jsonify(result), 200
        except Exception as e:
            return jsonify({'error': str(e)}), 500
        
       
  
    @app.route('/api/family_member_counts', methods=['GET'])
    def family_member_counts():
        """
        Returns a bar-chart payload of family-member counts per country.

        Extra logic:
        1️⃣ Read the most recent search keywords (same SQL used by
           /api/last_search_keywords).
        2️⃣ If none of those keywords appear in any patent title stored
           in raw_patents, we call the EPO-OPS updater
           (update_family_members_ops) to refresh the DB.
        3️⃣ After the optional refresh, aggregate family_jurisdictions
           exactly as before and shape the JSON for Chart.js.
        """
        try:
    # ───── 1️⃣  Read the column we care about ──────────────────────────
            query = text(
    'SELECT family_jurisdictions FROM raw_patents '
    'WHERE family_jurisdictions IS NOT NULL '
    '  AND array_length(family_jurisdictions, 1) > 0'
)
            df = pd.read_sql(query, engine)

    # ───── 2️⃣  If nothing there, run the heavy updater once ───────────
            if df.empty:
                app.logger.info("family_jurisdictions empty → running updater")
                update_family_members_ops()
                df = pd.read_sql(query, engine)      # retry

    # ───── 3️⃣  Still empty? Just return blanks ────────────────────────
            if df.empty:
                return jsonify({"labels": [], "datasets": []}), 200

    # ───── 4️⃣  Aggregate for the bar-chart ────────────────────────────
            df['country_codes'] = df['family_jurisdictions'].apply(
                lambda x: ast.literal_eval(x) if isinstance(x, str) else x
                    )

            from family_analysis import get_family_counts_by_country
            counts_df = get_family_counts_by_country(df[['country_codes']])

            return jsonify({
        "labels":   counts_df['country_code'].tolist(),
        "datasets": [{
            "label": "Family Member Count",
            "data":  counts_df['member_count'].tolist()
        }]
    }), 200

        except Exception as e:
            app.logger.error(f"Error in /api/family_member_counts: {e}")
            return jsonify({"error": str(e)}), 500
        
        
        
        
    @app.route('/api/family_size_distribution', methods=['GET'])
    def family_size_distribution():
        """
        Return a histogram-ready payload for the distribution of patent-family sizes
        stored in raw_patents.

        JSON shape
        ----------
        {
            "labels":   [1, 2, 3, ...],      # family size buckets
            "datasets": [{
                "label": "Number of Families",
                "data":  [432, 97, 26, ...]  # counts per bucket (same length)
            }]
        }
        """
        try:
            # ── 1️⃣  Grab non-empty family_members arrays ───────────────────
            query = text(
                'SELECT family_members FROM raw_patents '
                'WHERE family_members IS NOT NULL '
                '  AND array_length(family_members, 1) > 0'
            )
            df = pd.read_sql(query, engine)

            # ── 2️⃣  If DB has none, populate via the heavy updater ─────────
            if df.empty:
                app.logger.info("family_members empty → running updater")
                update_family_members_ops()
                df = pd.read_sql(query, engine)

            # ── 3️⃣  Still empty? return blanks so UI can handle gracefully ─
            if df.empty:
                return jsonify({"labels": [], "datasets": []}), 200

            # ── 4️⃣  Compute family_size per row and build histogram data ───
            def _to_list(x):
                # handle TEXT-stored arrays like '[US,EP]'
                return ast.literal_eval(x) if isinstance(x, str) else x

            df['family_size'] = df['family_members'].apply(_to_list).apply(len)

            counts = (
                df['family_size']
                .value_counts()
                .sort_index()       # ascending bucket order
            )

            return jsonify({
                "labels":   counts.index.tolist(),   # [1,2,3,...]
                "datasets": [{
                    "label": "Number of Families",
                    "data":  counts.values.tolist()
                }]
            }), 200

        except Exception as e:
            app.logger.error(f"Error in /api/family_size_distribution: {e}")
            return jsonify({"error": str(e)}), 500
        
        
        
        
        
        
        
        

    @app.route('/api/international_protection_matrix', methods=['GET'])
    def international_protection_matrix():
        """
        Returns a JSON payload describing, for each origin country
        (first-publication country), how many family members were filed
        in every other jurisdiction.

        Response shape
        --------------
        {
            "origins": ["US", "CN", "JP", ...],          # y-axis order (bottom → top)
            "filings": ["US", "CN", "EP", ...],          # x-axis order (largest → smallest)
            "matrix":  [                                 # counts[y][x]
        [123, 45,  9, ...],     # US row
        [ 67, 98, 10, ...],     # CN row
        ...
      ]
        }
    """
        try:
            # ── 1️⃣  Pull the two columns we need ───────────────────────────
            query = text(
            'SELECT first_publication_country AS origin_country, '
            '       family_jurisdictions                       '
            'FROM   raw_patents                                '
            'WHERE  first_publication_country IS NOT NULL      '
            '  AND  family_jurisdictions IS NOT NULL           '
            '  AND  array_length(family_jurisdictions, 1) > 0'
            )

            df = pd.read_sql(query, engine)

        # ── 2️⃣  If empty, run the heavy updater once & retry ───────────
            if df.empty:
                app.logger.info("family_jurisdictions empty → running updater")
                update_family_members_ops()
                df = pd.read_sql(query, engine)

            if df.empty:
                return jsonify({"origins": [], "filings": [], "matrix": []}), 200

        # ── 3️⃣  Build relationships origin → filing ─────────────────────
            df['family_list'] = df['family_jurisdictions'].apply(
                lambda x: ast.literal_eval(x) if isinstance(x, str) else x
            )

            records = []
            for idx, row in df.iterrows():
                for filing in row['family_list']:
                    records.append({
                        'origin_country': row['origin_country'],
                        'filing_country': filing
                    })

            rel_df = pd.DataFrame(records)

            grouped = (
                rel_df
                .groupby(['origin_country', 'filing_country'])
                .size()
                .reset_index(name='count')
            )

            pivot = (
                grouped
                .pivot(index='origin_country', columns='filing_country', values='count')
                .fillna(0)
            )

        # ── 4️⃣  Sort filing columns by total desc & origin rows by total asc
            pivot = pivot[pivot.sum(axis=0).sort_values(ascending=False).index]
            pivot = pivot.loc[pivot.sum(axis=1).sort_values(ascending=True).index]

            origins = pivot.index.tolist()
            filings = pivot.columns.tolist()
            matrix  = pivot.astype(int).values.tolist()

            return jsonify({
                "origins": origins,
                "filings": filings,
                "matrix":  matrix
            }), 200

        except Exception as e:
            app.logger.error(f"Error in /api/international_protection_matrix: {e}")
            return jsonify({"error": str(e)}), 500
        
        
        
        
        
        #  >>> NEW ────────────────────────────────────────────────────────────────
    @app.route('/api/international_patent_flow', methods=['GET'])
    def international_patent_flow():
        """
        Deliver a receiver-vs-origin matrix for the stacked horizontal
        “International Patent Flow” bar-chart.

        JSON shape
        ----------
        {
            "receivers": ["US", "CN", "EP", ...],     # y-axis order (top→bottom)
            "origins":   ["US", "CN", "JP", ...],     # legend / stack order
            "matrix":    [                            # counts[row][col]
                [731, 212,  65, ...],   # patents each receiver got from US, CN, JP…
                [452, 530, 110, ...],   # row for CN, etc.
                ...
            ]
        }
    """
        try:
            # ── 1️⃣  Pull origin + family_jurisdictions ────────────────────
            query = text(
                'SELECT first_publication_country AS origin_country, '
                '       family_jurisdictions                       '
                'FROM   raw_patents                                '
                'WHERE  first_publication_country IS NOT NULL      '
                '  AND  family_jurisdictions IS NOT NULL           '
                '  AND  array_length(family_jurisdictions, 1) > 0'
            )
            df = pd.read_sql(query, engine)

        # ── 2️⃣  Auto-populate if column empty ──────────────────────────
            if df.empty:
                app.logger.info("family_jurisdictions empty → running updater")
                update_family_members_ops()
                df = pd.read_sql(query, engine)

            if df.empty:
                return jsonify({"receivers": [], "origins": [], "matrix": []}), 200

        # ── 3️⃣  Build long-form receiver←origin relationships ─────────
            df['family_list'] = df['family_jurisdictions'].apply(
                lambda x: ast.literal_eval(x) if isinstance(x, str) else x
            )

            rel_records = []
            for idx, row in df.iterrows():
                for receiver in row['family_list']:
                    rel_records.append({
                        'receiver_country': receiver,
                        'origin_country':   row['origin_country']
                    })

            rel_df = pd.DataFrame(rel_records)

        # ── 4️⃣  Aggregate & pivot  receiver × origin  ─────────────────
            grouped = (
                rel_df
                .groupby(['receiver_country', 'origin_country'])
                .size()
                .reset_index(name='count')
            )

            pivot = (
                grouped
                .pivot(index='receiver_country', columns='origin_country', values='count')
                .fillna(0)
            )

        # ── 5️⃣  Row order: receivers sorted by total (desc)  ──────────
            pivot = pivot.loc[pivot.sum(axis=1).sort_values(ascending=False).index]

        #      Column order: origins sorted by overall total (desc) so
        #      biggest source stacks appear first (bottom of each bar)
            pivot = pivot[pivot.sum(axis=0).sort_values(ascending=False).index]

            receivers = pivot.index.tolist()
            origins   = pivot.columns.tolist()
            matrix    = pivot.astype(int).values.tolist()

            return jsonify({
                "receivers": receivers,
                "origins":   origins,
                "matrix":    matrix
            }), 200

        except Exception as e:
            app.logger.error(f"Error in /api/international_patent_flow: {e}")
            return jsonify({"error": str(e)}), 500
        
        
        
    @app.route('/api/geographical_distribution', methods=['GET'])
    def geographical_distribution():
        """
        Returns the patent count per first_publication_country.

        Example (top-3):
            {
              "labels": ["US", "JP", "CN"],
              "datasets": [{
                  "label": "Patent Count",
                  "data": [120, 83, 57],
                  "backgroundColor": "rgba(54, 162, 235, 0.7)"
              }]
            }
        """
        try:
            # 1️⃣  Pull the country column from raw_patents
            patents = RawPatent.query.with_entities(
                RawPatent.first_publication_country
            ).all()                     # -> list of tuples [(“US”,), (“JP”,)…]

            import pandas as pd
            df = pd.DataFrame(
                [c for (c,) in patents if c]  # drop NULL / empty
                , columns=['country']
            )

            if df.empty:                       # defensive – avoid div/0 on fresh DB
                return jsonify({
                    "labels": [],
                    "datasets": [{"label": "Patent Count", "data": []}]
                }), 200

            # 2️⃣  Aggregate & format for the card
            counts = df['country'].value_counts()
            result = {
                "labels": counts.index.tolist(),
                "datasets": [{
                    "label": "Patent Count",
                    "data": counts.values.tolist(),
                    "backgroundColor": "rgba(54, 162, 235, 0.7)"
                }]
            }
            return jsonify(result), 200

        except Exception as e:                 # ❸ standardised error handling
            app.logger.error(f"Error in geographical_distribution: {e}")
            return jsonify({"error": str(e)}), 500




    @app.route('/api/research_publications_by_year', methods=['GET'])
    def research_publications_by_year():
        """
        Retu1rns a list of {year, count} for publications in research_data3 table,
        suitable for a publication trend chart.
        """
        try:
            # Query all years from research_data3
            results = (
                db.session.query(ResearchData3.year, db.func.count(ResearchData3.id))
                .group_by(ResearchData3.year)
                .order_by(ResearchData3.year)
                .all()
            )
            # Build response
            data = [
                {"year": int(year), "count": int(count)}
                for year, count in results if year is not None
            ]
            return jsonify(data), 200
        except Exception as e:
            import traceback

            print(traceback.format_exc())
            return jsonify({"error": str(e)}), 500
      
      
      
      
        
    @app.route('/api/research/processed_texts', methods=['GET'])
    def research_processed_texts():
        """
        Research analogue of /api/processed_texts.
        Returns preprocessed strings built from title+abstract.
        """
        try:
            df = build_research_keyword_df()
            if df.empty:
                return jsonify({"error": "No records found in research_data3"}), 404

            processed_list = [txt for txt in df['processed_title'].tolist() if txt]
            if not processed_list:
                return jsonify({"error": "No valid text to process"}), 400

            return jsonify({"processed": processed_list}), 200
        except Exception as e:
            app.logger.error(f"Unexpected error in /api/research/processed_texts: {e}")
            return jsonify({"error": str(e)}), 500    
        
        
    @app.route('/api/research/topic_evolution', methods=['GET'])
    def research_topic_evolution():
        """
        Research analogue of /api/topic_evolution (patents).
        - Builds DF from research_data3 (title+abstract, year)
        - Calls research_analyze_topic_evolution (no fallback in route)
        - Stores ResearchWindow, ResearchDivergence, and ResearchTopic (mirrors patents)
        """
        try:
            # 1) Build DF and pre-clean
            keyword_df = build_research_keyword_df()
            keyword_df = keyword_df.dropna(subset=['title', 'year'])
            if keyword_df.empty:
                return jsonify({"message": "No research docs with valid title/year"}), 404

            # Avoid duplicate 'title' columns (keep preprocessed as title)
            keyword_df = keyword_df.drop(columns=['title']).rename(columns={'processed_title': 'title'})
            keyword_df['title'] = keyword_df['title'].astype(str)

            # Optional min_year (?min_year=YYYY)
            req_min_year = request.args.get('min_year', default=None, type=int)
            if req_min_year is not None:
                yint = pd.to_numeric(keyword_df['year'], errors='coerce').fillna(-1).astype(int)
                mask = (yint >= req_min_year).to_numpy()
                keyword_df = keyword_df.loc[mask].copy()

            if keyword_df.empty:
                return jsonify({"topic_evolution": [], "windows": [], "divergences": []}), 200

            # 2) Analyze topic evolution (no route-level fallback)
            topic_evolution_data, windows_data, divergences, valid_years = research_analyze_topic_evolution(keyword_df)

            # 3) Clear and store ResearchWindow + ResearchDivergence + ResearchTopic
            db.session.query(ResearchTopic).delete()
            db.session.query(ResearchWindow).delete()
            db.session.query(ResearchDivergence).delete()
            db.session.commit()

            # Windows
            window_objects = []
            for win in windows_data:
                w = ResearchWindow(start_year=int(win['start']), end_year=int(win['end']))
                db.session.add(w)
                window_objects.append(w)
            db.session.commit()

            # Map (start, end) → window obj
            window_idx = {(w.start_year, w.end_year): w for w in window_objects}

            # Topics (store words+weights compactly in JSON string)
            for t in topic_evolution_data:
                w = window_idx.get((int(t['start']), int(t['end'])))
                if not w:
                    continue
                topic_id_str  = t['topic_id']
                topic_number  = int(topic_id_str.split('-')[-1])
                words         = t['words']
                weights       = [float(x) for x in t['weights']]
                # short description: top 5 terms with weights
                desc = "; ".join(f"{w_} ({weights[i]:.2f})" for i, w_ in enumerate(words[:5]))
                fields_json = json.dumps({"words": words, "weights": weights})
                db.session.add(ResearchTopic(
                    window_id=w.id,
                    topic=str(topic_number),
                    description=desc,
                    fields=fields_json
                ))
            db.session.commit()

            # Divergences
            for i in range(len(divergences)):
                db.session.add(ResearchDivergence(
                    from_year=int(valid_years[i]),
                    to_year=int(valid_years[i + 1]),
                    divergence=float(divergences[i])
                ))
            db.session.commit()

            # 4) Query back & shape response (like patents endpoint)
            windows = ResearchWindow.query.order_by(ResearchWindow.start_year).all()
            topics  = (db.session.query(ResearchTopic)
                       .join(ResearchWindow, ResearchTopic.window_id == ResearchWindow.id)
                       .order_by(ResearchWindow.start_year, ResearchTopic.topic)
                       .all())
            divergences_q = ResearchDivergence.query.order_by(ResearchDivergence.from_year).all()

            window_list = [{
                'start': w.start_year,
                'end':   w.end_year,
                'years': list(range(w.start_year, w.end_year + 1))
            } for w in windows]

            topic_evolution_list = []
            for t in topics:
                w = t.window
                data = json.loads(t.fields or '{}')
                topic_evolution_list.append({
                    'start':  w.start_year,
                    'end':    w.end_year,
                    'topic_id': f"{w.start_year}-{int(t.topic)}",
                    'words':   data.get('words', []),
                    'weights': data.get('weights', [])
                })

            divergences_list = [{
                'from_year': d.from_year,
                'to_year':   d.to_year,
                'divergence': d.divergence
            } for d in divergences_q]

            return jsonify({
                "topic_evolution": topic_evolution_list,
                "windows": window_list,
                "divergences": divergences_list
            }), 200

        except Exception as e:
            db.session.rollback()
            app.logger.error(f"Error in /api/research/topic_evolution: {e}")
            return jsonify({"error": f"Failed to process research topic evolution: {str(e)}"}), 500
        
        



    @app.route('/api/research/weighted_word_clouds', methods=['GET'])
    def research_weighted_word_clouds():
        """
        Research analogue of /api/weighted_word_clouds (patents).
        - Reads windows from ResearchWindow
        - Reads topics from ResearchTopic; parses topic.fields JSON: {"words": [...], "weights": [...]} 
        - Aggregates weights per word across all topics within the same window
        - Returns: [{ start, end, words: [{word, weight}, ...] }, ...]
        """
        try:
            logger = logging.getLogger(__name__)

            # 1) Fetch windows in order
            windows = ResearchWindow.query.order_by(ResearchWindow.start_year).all()
            result = []

            for window in windows:
                # dict to sum weights per normalized word
                word_weight_dict = defaultdict(float)

                # 2) Get all topics belonging to this window
                topics = ResearchTopic.query.filter_by(window_id=window.id).all()

                for topic in topics:
                    # 2.a) Parse the fields JSON safely
                    try:
                        payload = json.loads(topic.fields or "{}")
                        words   = payload.get("words", []) or []
                        weights = payload.get("weights", []) or []
                    except Exception as e:
                        logger.warning(f"Skipping research topic {topic.id}: invalid JSON in fields. err={e}")
                        continue

                    # 2.b) Validate aligned lengths
                    if len(words) != len(weights):
                        logger.warning(
                            f"Skipping research topic {topic.id}: words/weights length mismatch "
                            f"({len(words)} vs {len(weights)})"
                        )
                        continue

                    # 2.c) Aggregate weights per word (normalize token, but keep original text as key)
                    for w, wt in zip(words, weights):
                        # normalize like patents endpoint
                        processed = preprocess_text((w or "").lower())
                        if processed and processed.strip():
                            try:
                                word_weight_dict[w] += float(wt)
                            except Exception:
                                # ignore non-numeric weights silently
                                continue

                    if not word_weight_dict:
                        logger.info(f"No valid research words for window {window.start_year}-{window.end_year}")
                        continue

                    # 3) Sort words by aggregated weight desc
                    word_list = [
                        {"word": w, "weight": weight}
                        for w, weight in sorted(word_weight_dict.items(), key=lambda x: x[1], reverse=True)
                    ]

                # 4) Append this window’s cloud
                result.append({
                    "start": window.start_year,
                    "end": window.end_year,
                    "words": word_list
                })

            logger.info(f"Generated research word cloud data for {len(result)} windows")
            return jsonify(result), 200

        except Exception as e:
            logger.error(f"Error in research_weighted_word_clouds endpoint: {str(e)}")
            return jsonify({"error": str(e)}), 500

    
    
    
    @app.route('/api/research/automatic_topic_shift', methods=['GET'])
    def research_automatic_topic_shift():
        """
        Research analogue of /api/automatic_topic_shift.

        - Reads divergence deltas between consecutive years from research_divergences
        - Computes a robust threshold (80th percentile) to flag "large" shifts
        - Returns both the raw divergences and the window boundaries from research_windows
        """
        try:
            # 1) Pull divergence series for research papers (ordered by year)
            divergences = ResearchDivergence.query.order_by(ResearchDivergence.from_year).all()
            if not divergences:
                # Same behavior as patents: return 200 with a friendly message if nothing to show
                return jsonify({"message": "No research divergence data available"}), 200

            # 2) Threshold to highlight big jumps (80th percentile is a good default)
            divergence_values = [float(d.divergence) for d in divergences]
            threshold = np.percentile(divergence_values, 80) if divergence_values else 0.0

            # 3) Pull detected time windows for the research corpus
            windows = ResearchWindow.query.order_by(ResearchWindow.start_year).all()

            # 4) Format payloads for the UI
            divergence_data = [
                {
                    "from_year": int(d.from_year),
                    "to_year": int(d.to_year),
                    "divergence": float(d.divergence),
                }
                for d in divergences
            ]

            window_list = [
                {
                    "start": int(w.start_year),
                    "end": int(w.end_year),
                }
                for w in windows
            ]

            # 5) Ship JSON (mirrors patent endpoint structure for easy frontend reuse)
            response = {
                "divergence_data": divergence_data,
                "threshold": float(threshold),
                "windows": window_list,
            }
            return jsonify(response), 200

        except Exception as e:
            # Mirror the patent endpoint’s error handling
            return jsonify({"error": str(e)}), 500
    
        
    
    
    @app.route('/api/research/cooccurrence_trends', methods=['GET'])
    def research_cooccurrence_trends():
        """
        Returns JSON data for plotting co-occurrence trends of *research* keyword pairs.
        - Uses research_data3 (title + abstract + year)
        - Removes English & French stopwords
        - Ensures term pairs are not identical (case-insensitive)
        Query params (optional):
        - window_size: sliding window size for co-occurrence (default 5)
        - min_count:   minimum total co-occurrence count to keep a pair (default 10)
        - top:         number of top emerging/declining pairs to include (default 5)
        """
        try:
            # ---- 0) Read params (defaults mirror the patent endpoint) ----
            window_size = int(request.args.get('window_size', 5))
            min_count   = int(request.args.get('min_count', 10))
            top_k       = int(request.args.get('top', 5))

            # ---- 1) Fetch research data (title + abstract + year) ----
            # We concatenate title & abstract to enrich context before co-occurrence.
            query = 'SELECT title, abstract, year FROM research_data3'
            df = pd.read_sql(query, engine)

            # Keep only rows with a valid year
            df['year'] = pd.to_numeric(df['year'], errors='coerce').astype('Int64')
            df = df.dropna(subset=['year'])
            df['year'] = df['year'].astype(int)

            # ---- 2) Build unified text field and clean stopwords ----
            df['text'] = (df['title'].fillna('') + ' ' + df['abstract'].fillna('')).str.strip()
            # remove EN+FR stopwords & punctuation (same helper as the patents endpoint)
            df['text'] = df['text'].apply(clean_text_remove_stopwords)  # :contentReference[oaicite:1]{index=1}

            # Remove empty docs after cleaning
            df = df[df['text'].str.strip().astype(bool)]
            if df.empty:
                return jsonify({
                    "emerging": [],
                    "declining": [],
                    "emerging_pairs": [],
                    "declining_pairs": []
                }), 200

            # ---- 3) Prepare grouped DataFrame (one concatenated string per year) ----
            # This mirrors what you do for patents: group texts within the same year, then join.
            grouped = df.groupby('year')['text'].apply(lambda texts: " ".join(texts)).reset_index()

            # ---- 4) Run co-occurrence trend analysis ----
            # Sliding-window co-occurrence + per-year frequencies + linear trend (slope, p-value).
            cooc_trends = track_cooccurrence_trends(
                grouped,
                time_col='year',
                text_col='text',
                window_size=window_size,
                min_count=min_count
            )  # :contentReference[oaicite:2]{index=2}

            if cooc_trends.empty:
                return jsonify({
                    "emerging": [],
                    "declining": [],
                    "emerging_pairs": [],
                    "declining_pairs": []
                }), 200

            # ---- 5) Filter out identical pairs (case-insensitive) ----
            cooc_trends = cooc_trends[
                cooc_trends['term1'].str.lower() != cooc_trends['term2'].str.lower()
            ]

            # ---- 6) Select top emerging & declining pairs (significant slopes) ----
            emerging_tech = cooc_trends[
                (cooc_trends.slope > 0) & (cooc_trends.p_value < 0.05)
            ].sort_values('slope', ascending=False).head(top_k)

            declining_tech = cooc_trends[
                (cooc_trends.slope < 0) & (cooc_trends.p_value < 0.05)
            ].sort_values('slope').head(top_k)

            # ---- 7) Shape time-series points for plotting ----
            def prepare_plot_data(df_in):
                plot_data = []
                for _, row in df_in.iterrows():
                    # frequency_history is a list of (year, freq) pairs
                    for year, freq in row['frequency_history']:
                        plot_data.append({
                            'year': int(year),
                            'frequency': int(freq),
                            'term_pair': f"{row['term1']} & {row['term2']}"
                        })
                return plot_data

            response = {
                "emerging": prepare_plot_data(emerging_tech),
                "declining": prepare_plot_data(declining_tech),
                "emerging_pairs": [
                    {
                        "term1": row['term1'],
                        "term2": row['term2'],
                        "slope": row['slope'],
                        "total_count": row['total_count']
                    }
                    for _, row in emerging_tech.iterrows()
                ],
                "declining_pairs": [
                    {
                        "term1": row['term1'],
                        "term2": row['term2'],
                        "slope": row['slope'],
                        "total_count": row['total_count']
                    }
                    for _, row in declining_tech.iterrows()
                ]
            }
            return jsonify(response), 200

        except Exception as e:
            return jsonify({"error": str(e)}), 500

    
    
    @app.route('/api/research/citations_per_year_stats', methods=['GET'])
    def research_citations_per_year_stats():
        """
        Analysis #1 — Age-normalized impact per cohort (citations per year).
        Uses research_data3.year and research_data3.citation_count.
        Optional query params:
        - min_year, max_year: restrict the cohort range.
        Returns a list of {year, n_papers, total_citations, mean_citations, median_citations,
                           mean_citations_per_year, median_citations_per_year}.
        """
        try:
            # 1) Pull minimal columns from DB
            rows = (
                db.session.query(ResearchData3.year, ResearchData3.citation_count)
                .filter(ResearchData3.year.isnot(None))
                .all()
            )
            if not rows:
                return jsonify([]), 200

            # 2) To DataFrame and drop NaNs
            df = pd.DataFrame(rows, columns=['year', 'citation_count']).dropna()
            # Enforce integer year and non-negative citations
            df['year'] = pd.to_numeric(df['year'], errors='coerce').astype('Int64')
            df['citation_count'] = pd.to_numeric(df['citation_count'], errors='coerce').fillna(0).clip(lower=0).astype(int)
            df = df.dropna(subset=['year'])
            df['year'] = df['year'].astype(int)

            # Optional year filter
            min_year = request.args.get('min_year', type=int)
            max_year = request.args.get('max_year', type=int)
            if min_year is not None:
                df = df[df['year'] >= min_year]
            if max_year is not None:
                df = df[df['year'] <= max_year]
            if df.empty:
                return jsonify([]), 200

            # 3) Age-normalize citations: CPY = citations / (current_year - year + 1)
            current_year = datetime.now().year
            age = (current_year - df['year'] + 1).clip(lower=1)
            df = df.assign(cpy=df['citation_count'] / age)

            # 4) Aggregate per year
            grouped = df.groupby('year')
            out = []
            for yr, g in grouped:
                cits = g['citation_count'].to_numpy()
                cpy = g['cpy'].to_numpy()
                out.append({
                    "year": int(yr),
                    "n_papers": int(len(g)),
                    "total_citations": int(cits.sum()),
                    "mean_citations": float(cits.mean()),
                    "median_citations": float(np.median(cits)),
                    "mean_citations_per_year": float(cpy.mean()),
                    "median_citations_per_year": float(np.median(cpy)),
                })
            out.sort(key=lambda d: d["year"])
            return jsonify(out), 200

        except Exception as e:
            return jsonify({"error": str(e)}), 500

    
    
    @app.route('/api/research/citation_percentiles', methods=['GET'])
    def research_citation_percentiles():
        """
        Analysis #2 — Percentiles & top-decile concentration per year.
        For each year: P50, P90, P99 of citation_count, plus the share of total citations
        captured by the top 10% most-cited papers (by that year's distribution).
        Optional query params: min_year, max_year
        """
        try:
            rows = (
                db.session.query(ResearchData3.year, ResearchData3.citation_count)
                .filter(ResearchData3.year.isnot(None))
                .all()
            )
            if not rows:
                return jsonify([]), 200

            df = pd.DataFrame(rows, columns=['year', 'citation_count']).dropna()
            df['year'] = pd.to_numeric(df['year'], errors='coerce').astype('Int64')
            df['citation_count'] = pd.to_numeric(df['citation_count'], errors='coerce').fillna(0).clip(lower=0).astype(int)
            df = df.dropna(subset=['year'])
            df['year'] = df['year'].astype(int)

            min_year = request.args.get('min_year', type=int)
            max_year = request.args.get('max_year', type=int)
            if min_year is not None:
                df = df[df['year'] >= min_year]
            if max_year is not None:
                df = df[df['year'] <= max_year]
            if df.empty:
                return jsonify([]), 200

            out = []
            for yr, g in df.groupby('year'):
                arr = g['citation_count'].to_numpy()
                if arr.size == 0:
                    continue
                total = arr.sum()
                p50 = float(np.percentile(arr, 50))
                p90 = float(np.percentile(arr, 90))
                p99 = float(np.percentile(arr, 99))
                # top-decile threshold and what fraction of total citations it captures
                thresh = np.percentile(arr, 90)
                top_mask = arr >= thresh
                top_share = float(arr[top_mask].sum() / total) if total > 0 else 0.0
                out.append({
                    "year": int(yr),
                    "n_papers": int(arr.size),
                    "p50": p50,
                    "p90": p90,
                    "p99": p99,
                    "top10_threshold": float(thresh),
                    "top10_share_of_citations": top_share
                })
            out.sort(key=lambda d: d["year"])
            return jsonify(out), 200

        except Exception as e:
            return jsonify({"error": str(e)}), 500

    
    
    
        
    def gini_coefficient(x: np.ndarray) -> float:
        """Compute Gini on non-negative vector x. Returns 0..1. Empty or all-zero → 0."""
        x = np.asarray(x, dtype=float)
        if x.size == 0:
            return 0.0
        if np.all(x == 0):
            return 0.0
        if np.any(x < 0):
            x = np.clip(x, 0, None)
        # Mean absolute difference / (2 * mean)
        mad = np.abs(np.subtract.outer(x, x)).mean()
        mean = x.mean()
        return float(mad / (2 * mean))


    @app.route('/api/research/citation_inequality', methods=['GET'])
    def research_citation_inequality():
        """
        Analysis #3 — Inequality of attention.
        For each year: Gini coefficient of citations and concentration captured by top decile & top 1%.
        Optional query params: min_year, max_year
        """
        try:
            rows = (
                db.session.query(ResearchData3.year, ResearchData3.citation_count)
                .filter(ResearchData3.year.isnot(None))
                .all()
            )
            if not rows:
                return jsonify([]), 200

            df = pd.DataFrame(rows, columns=['year', 'citation_count']).dropna()
            df['year'] = pd.to_numeric(df['year'], errors='coerce').astype('Int64')
            df['citation_count'] = pd.to_numeric(df['citation_count'], errors='coerce').fillna(0).clip(lower=0).astype(int)
            df = df.dropna(subset=['year'])
            df['year'] = df['year'].astype(int)

            min_year = request.args.get('min_year', type=int)
            max_year = request.args.get('max_year', type=int)
            if min_year is not None:
                df = df[df['year'] >= min_year]
            if max_year is not None:
                df = df[df['year'] <= max_year]
            if df.empty:
                return jsonify([]), 200

            out = []
            for yr, g in df.groupby('year'):
                arr = g['citation_count'].to_numpy()
                total = arr.sum()
                gini = gini_coefficient(arr)
                # Decile & top-1% concentration
                d90 = np.percentile(arr, 90)
                d99 = np.percentile(arr, 99)
                top10_share = float(arr[arr >= d90].sum() / total) if total > 0 else 0.0
                top1_share  = float(arr[arr >= d99].sum() / total) if total > 0 else 0.0

                out.append({
                    "year": int(yr),
                    "n_papers": int(arr.size),
                    "gini": gini,
                    "top10_share_of_citations": top10_share,
                    "top1_share_of_citations": top1_share
                })
            out.sort(key=lambda d: d["year"])
            return jsonify(out), 200

        except Exception as e:
            return jsonify({"error": str(e)}), 500

        
        
    @app.route('/api/research/papers', methods=['GET'])
    def get_research_papers():
        """
        Fetches paper rows from research_data3 for the Papers tab.
        No search parsing here—just returns stored results after /api/scientific_search_merge has populated the table.
        Query params (optional):
        - limit: max rows to return (default 500, max 2000)
        - sort_by: 'citations' (default) or 'year'
        - order: 'desc' (default) or 'asc'
        Returns: [{ title, year, journal_name, fields_of_study, citation_count, reference_count, subcategory }, ...]
        """
        try:
            # --- read query params with guardrails ---
            limit = int(request.args.get('limit', 500))
            limit = max(1, min(limit, 2000))
            sort_by = (request.args.get('sort_by') or 'citations').lower()
            order = (request.args.get('order') or 'desc').lower()

            # --- choose sort column ---
            order_col = ResearchData3.citation_count if sort_by != 'year' else ResearchData3.year
            if order == 'asc':
                order_by = order_col.asc().nullslast()
            else:
                order_by = order_col.desc().nullslast()

            # --- query rows ---
            rows = (
                db.session.query(ResearchData3)
                .order_by(order_by)
                .limit(limit)
                .all()
        )

            # --- serializer: keep exactly the fields your table needs ---
            def parse_fos(fos_raw):
                # fields_of_study may be JSON string, plain string, or list
                if isinstance(fos_raw, list):
                    return fos_raw
                if isinstance(fos_raw, str):
                    try:
                        tmp = json.loads(fos_raw)
                        if isinstance(tmp, list):
                            return tmp
                        if isinstance(tmp, str):
                            return [tmp]
                    except Exception:
                        # comma-separated plain text
                        return [s.strip() for s in fos_raw.split(",") if s.strip()]
                return None

            out = []
            for r in rows:
                journal = getattr(r, "journal_name", None) or getattr(r, "publication_venue_name", None)
                out.append({
                    "title": getattr(r, "title", None),
                    "year": getattr(r, "year", None),
                    "journal_name": journal,
                    "fields_of_study": parse_fos(getattr(r, "fields_of_study", None)),
                    "citation_count": getattr(r, "citation_count", None),
                    "reference_count": getattr(r, "reference_count", None),
                    "subcategory": getattr(r, "subcategory", None),
                })

            return jsonify(out), 200

        except Exception as e:
            app.logger.error(f"/api/research/papers error: {e}")
            return jsonify({"error": str(e)}), 500
  
    def _year_to_ds(y: int) -> str:
    # ISO-like date at year end (Prophet uses Dec-31)
        return f"{int(y)}-12-31"

    def _safe_float(x):
        try:
            if x is None or (isinstance(x, float) and np.isnan(x)):
                return None
            return float(x)
        except Exception:
            return None
        
        
        
        
        
        



    @app.route("/api/prophet_forecast", methods=["GET", "POST"])
    def api_prophet_forecast():
        """
        Prophet single-series forecasts + growth-rate diagnostics.

        NEW in this version:
          - Computes the "past" growth rate twice:
              (A) Actual past GR from historical counts (unchanged).
              (B) Counterfactual past GR using predictions only:
                  Train on data up to the year BEFORE the past-window,
                  forecast the whole past-window years, then compute GR
              over that window using the predicted counts. This lets
              you compare model-vs-actual on the exact historical window.
        """
        try:
            payload = request.get_json(silent=True) or {}

            # ------------------ Parse inputs (unchanged) ------------------
            horizon    = int(payload.get("horizon", HORIZON_DEFAULT))
            horizon    = validate_horizon(horizon)
            test_years = int(payload.get("test_years", request.args.get("test_years", TEST_YEARS_DEFAULT)))
            pub_tail   = int(payload.get("pub_tail",   request.args.get("pub_tail",   PUB_TAIL_TRUNC_DEFAULT)))
            pat_tail   = int(payload.get("pat_tail",   request.args.get("pat_tail",   PAT_TAIL_TRUNC_DEFAULT)))
            max_lag    = int(payload.get("max_lag",    request.args.get("max_lag",    MAX_LAG_DEFAULT)))
            tech_label =       payload.get("tech_label", request.args.get("tech_label", "current"))

            # Optional explicit split controls (kept for compatibility)
            split_year      = payload.get("split_year",      request.args.get("split_year"))
            eval_start_year = payload.get("eval_start_year", request.args.get("eval_start_year"))
            eval_end_year   = payload.get("eval_end_year",   request.args.get("eval_end_year"))
            split_year      = int(split_year)      if split_year is not None else None
            eval_start_year = int(eval_start_year) if eval_start_year is not None else None
            eval_end_year   = int(eval_end_year)   if eval_end_year is not None else None

            # ------------------ Run main pipeline (unchanged) ------------------
            pubs_fc, pats_fc, metrics_df, summary, pubs_test_eval, pats_test_eval = run_for_current_series(
                tech_label=tech_label,
                pub_tail_trunc=pub_tail,
                pat_tail_trunc=pat_tail,
                max_lag=max_lag,
                test_years=test_years,
                horizon=horizon,
                split_year=split_year,
                eval_start_year=eval_start_year,
                eval_end_year=eval_end_year,
            )

            # Keep history for plotting
            engine = _require_engine()
            pubs_hist, pats_hist = fetch_current_series(
                engine=engine,
                pub_tail_trunc=pub_tail,
                pat_tail_trunc=pat_tail
            )

            # Prepare history payloads (unchanged)
            pubs_history = [
                {"year": int(r.year), "ds": _year_to_ds(int(r.year)), "count": int(r.pub_count)}
                for r in pubs_hist.itertuples(index=False)
            ]
            pats_history = [
                {"year": int(r.year), "ds": _year_to_ds(int(r.year)), "count": int(r.patent_count)}
                for r in pats_hist.itertuples(index=False)
            ]

            # ---------- Build production forecast lists (unchanged) ----------
            pubs_forecast = [
                {
                    "year": int(r.year),
                    "ds": _year_to_ds(int(r.year)),
                    "yhat": _safe_float(r.pub_count_hat),
                    "yhat_lower": _safe_float(r.yhat_lower),
                    "yhat_upper": _safe_float(r.yhat_upper),
                }
                for r in pubs_fc.itertuples(index=False)
            ]
            pats_with = [
                {"year": int(r.year), "ds": _year_to_ds(int(r.year)), "yhat": _safe_float(r.yhat_with_pub_reg)}
                for r in pats_fc.itertuples(index=False)
            ]
            pats_base = [
                {"year": int(r.year), "ds": _year_to_ds(int(r.year)), "yhat": _safe_float(r.yhat_baseline)}
                for r in pats_fc.itertuples(index=False)
            ]

            # ---------- Use the already-computed TEST predictions (unchanged) ----------
            pubs_test_points = [
                {
                    "year": int(r.year),
                    "ds": _year_to_ds(int(r.year)),
                    "actual": _safe_float(r.actual),
                    "yhat": _safe_float(r.yhat),
                    "yhat_lower": _safe_float(r.yhat_lower),
                    "yhat_upper": _safe_float(r.yhat_upper),
                }
                for r in pubs_test_eval.itertuples(index=False)
            ]
            pats_test_points = [
                {
                    "year": int(r.year),
                    "ds": _year_to_ds(int(r.year)),
                    "actual": _safe_float(r.actual),
                    "yhat_baseline": _safe_float(r.yhat_baseline),
                    "yhat_with_pub_reg": _safe_float(r.yhat_with_pub_reg),
                }
                for r in pats_test_eval.itertuples(index=False)
            ]

            # ---------- Metrics (unchanged) ----------
            md = metrics_df.iloc[0].to_dict() if len(metrics_df) else {}
            metrics = {
                "mae_pubs": _safe_float(md.get("mae_pubs")),
                "rmse_pubs": _safe_float(md.get("rmse_pubs")),
                "ampe_pubs": _safe_float(md.get("ampe_pubs")),
                "mae_patents_baseline": _safe_float(md.get("mae_patents_baseline")),
                "rmse_patents_baseline": _safe_float(md.get("rmse_patents_baseline")),
                "ampe_patents_baseline": _safe_float(md.get("ampe_patents_baseline")),
                "mae_patents_with_reg": _safe_float(md.get("mae_patents_with_reg")),
                "rmse_patents_with_reg": _safe_float(md.get("rmse_patents_with_reg")),
                "ampe_patents_with_reg": _safe_float(md.get("ampe_patents_with_reg")),
            }

            # =====================================================================
            # GROWTH RATE BLOCKS
            # =====================================================================

            # 1) Build a combined "actual + future" counts table for main plotting
            #    (unchanged behavior for main dashboard).
            pats_actual = pats_hist.loc[:, ["year", "patent_count"]].copy()
            pats_future = pats_fc.loc[:, ["year", "yhat_with_pub_reg"]].rename(
                columns={"yhat_with_pub_reg": "patent_count"}
            )
            pats_counts_full = (
                pd.concat([pats_actual, pats_future], ignore_index=True)
                  .sort_values(["year"])
                  .drop_duplicates(subset=["year"], keep="last")
                  .sort_values("year")
                  .reset_index(drop=True)
            )

            # 2) Compute the "actual" past and current GR exactly as before.
            past_GR, past_start, past_end, past_label, curr_GR, curr_start, curr_end, curr_label = \
                compute_patent_growth_from_counts(pats_counts_full)

            # 3) NEW: Counterfactual "past" GR using forecasts only on the past window.
            #    We re-run the pipeline with the evaluation window set to [past_start..past_end],
            #    which ensures TRAIN ≤ (past_start-1) and TEST = that past window. We then
            #    stitch the counts as [actual up to past_start-1] + [predicted for past window]
            #    and compute GR again over the same window.
            past_forecast_GR = None
            past_forecast_label = None

            if past_start is not None and past_end is not None and past_start <= past_end:
                # Re-run with train cutoff at (past_start - 1), test on the past window
                _, _, _, _, _, pats_test_eval_past = run_for_current_series(
                    tech_label=tech_label,
                    pub_tail_trunc=pub_tail,
                    pat_tail_trunc=pat_tail,
                    max_lag=max_lag,
                    test_years=(past_end - past_start + 1),
                    horizon=horizon,  # not used for this eval
                    split_year=None,
                    eval_start_year=past_start,
                    eval_end_year=past_end,
                )

                if len(pats_test_eval_past):
                    # Build the counterfactual series:
                    #   - actual counts up to (past_start - 1)
                    #   - predicted counts for [past_start .. past_end]
                    before_mask = pats_hist["year"] < past_start
                    prior_actual = pats_hist.loc[before_mask, ["year", "patent_count"]].copy()

                    preds = pats_test_eval_past.loc[:, ["year", "yhat_with_pub_reg"]].rename(
                        columns={"yhat_with_pub_reg": "patent_count"}
                    )

                    pats_counterfactual = (
                        pd.concat([prior_actual, preds], ignore_index=True)
                        .sort_values("year")
                        .reset_index(drop=True)
                    )

                    # Compute GR again on this "actual+predicted-past" timeline.
                    past_cf_GR, cf_start, cf_end, past_cf_label, *_ = \
                        compute_patent_growth_from_counts(pats_counterfactual)

                    # Keep only the past-window values
                    past_forecast_GR = past_cf_GR
                    past_forecast_label = past_cf_label

            # Pack growth for response
            def _pack_growth(percent, label, start_year, end_year):
                return {
                    "percent": _safe_float(percent),
                    "label": label,
                    "window": {
                        "start_year": int(start_year) if start_year is not None else None,
                        "end_year": int(end_year) if end_year is not None else None,
                    },
                }

            growth = {
                # Unchanged: actual past and current on (actual + future) series
                "past": _pack_growth(past_GR, past_label, past_start, past_end),
                "current": _pack_growth(curr_GR, curr_label, curr_start, curr_end),
            }

            # New: model-vs-actual “past” comparison (if computed)
            if past_forecast_GR is not None:
                growth["past_counterfactual"] = _pack_growth(
                    past_forecast_GR, past_forecast_label, past_start, past_end
                )
                growth["past_compare"] = {
                    "delta_percent": _safe_float(past_forecast_GR - past_GR) if past_GR is not None else None
                }

            # ------------------ Shape final response (unchanged) ------------------
            resp = {
                "tech": summary.tech,
                "best_lag": int(summary.best_lag),
                "xcorr": _safe_float(summary.xcorr),
                "metrics": metrics,
                "growth": growth,
                "publications": {
                    "history": pubs_history,
                    "test": pubs_test_points,
                    "forecast": pubs_forecast,
                },
                "patents": {
                    "history": pats_history,
                    "test": pats_test_points,
                    "forecast": {
                        "with_pub_ar": pats_with,
                        "baseline": pats_base
                    }
                }
            }
        

            return jsonify(resp), 200

        except Exception as e:
            return jsonify({"error": str(e)}), 500




        


    EPS = 1e-9

    def _get_int(payload, name, default, minv=None, maxv=None):
        if name in payload:
            try: v = int(str(payload[name]).strip())
            except: v = default
        else:
            qv = request.args.get(name)
            if qv is None: v = default
            else:
                try: v = int(qv)
                except: v = default
        if minv is not None: v = max(minv, v)
        if maxv is not None: v = min(maxv, v)
        return v

    def _mspe(y_true, y_pred, eps=EPS):
        yt = np.asarray(y_true, float)
        yp = np.asarray(y_pred, float)
        m = np.isfinite(yt) & np.isfinite(yp)
        if not np.any(m): return float('nan')
        denom = np.maximum(np.abs(yt[m]), eps)
        return float(np.mean(((yt[m] - yp[m]) / denom) ** 2) * 100.0)

    def _build_enhanced(df_hist):
        df_enh = create_enhanced_features(df_hist.copy())
        feature_cols = [
            'patents','publications',
            'patents_ma3','publications_ma3',
            'patents_growth','publications_growth',
        'patent_pub_ratio','total_activity',
        'patents_lag1','publications_lag1',
        'patents_volatility','publications_volatility',
        'patents_trend','publications_trend'
        ]
        X = df_enh[feature_cols].astype(float).values
        years = df_enh['year'].astype(int).values
        return df_enh, X, years, feature_cols

    def _make_windows(X, years, n_steps):
        X_all, tgt_years = [], []
        for i in range(len(X) - n_steps):
            X_all.append(X[i:i+n_steps])
            tgt_years.append(int(years[i+n_steps]))
        return np.array(X_all), np.array(tgt_years, int)

    def _upsert_counts(df_counts, year, value):
        df = df_counts.copy()
        if (df['year'] == year).any():
            df.loc[df['year'] == year, 'patent_count'] = float(value)
        else:
            df = pd.concat([df, pd.DataFrame([{'year': int(year), 'patent_count': float(value)}])], ignore_index=True)
        return df.sort_values('year').reset_index(drop=True)

    @app.route('/api/lstm_forecast', methods=['GET', 'POST'])
    def lstm_forecast():
        try:
            payload    = request.get_json(silent=True) or {}
            horizon    = _get_int(payload, "horizon",   5, 1)
            pat_trunc  = _get_int(payload, "pat_trunc", 0, 0)
            pub_trunc  = _get_int(payload, "pub_trunc", 0, 0)
            test_years = _get_int(payload, "test_years", 0, 0)
            split_year = payload.get("split_year", request.args.get("split_year"))
            split_year = int(split_year) if split_year not in (None, "", "null") else None

            # ---- model + scalers (MUST match training) ----
            model_path  = os.path.join("backend/outputs", "lstm_patent_forecaster_enhanced.keras")
            scaler_path = os.path.join("backend/outputs", "global_scalers.pkl")
            model   = load_model(model_path)
            scalers = joblib.load(scaler_path)
            x_scaler = scalers["x_scaler"] if isinstance(scalers, dict) else scalers.x_scaler
            y_scaler = scalers["y_scaler"] if isinstance(scalers, dict) else scalers.y_scaler

            # ✅ infer step length from model (no hard-coded N_STEPS)
            n_steps = int((model.input_shape[1] if isinstance(model.input_shape, (list, tuple)) else 5) or 5)

            # ---- fetch DB data ----
            with engine.connect() as conn:
                df_pat_raw = pd.read_sql(text("""
                    SELECT first_filing_year AS year
                    FROM raw_patents
                    WHERE first_filing_year IS NOT NULL
                """), conn)
            if df_pat_raw.empty:
                return jsonify({"ok": False, "error": "No patent data found"}), 404

            with engine.connect() as conn:
                df_pub_raw = pd.read_sql(text("""
                    SELECT year
                FROM research_data3
                WHERE year IS NOT NULL
                """), conn)
            if df_pub_raw.empty:
                return jsonify({"ok": False, "error": "No publication data found"}), 404

            # ---- truncation ----
            max_pat_before = int(df_pat_raw["year"].max())
            if pat_trunc > 0:
                df_pat_raw = df_pat_raw[df_pat_raw["year"] <= max_pat_before - pat_trunc]
            if df_pat_raw.empty:
                return jsonify({"ok": False, "error": "No patent data left after truncation"}), 400
            max_pat_after = int(df_pat_raw["year"].max())
            eff_pat_trunc = max_pat_before - max_pat_after

            max_pub_before = int(df_pub_raw["year"].max())
            if pub_trunc > 0:
                df_pub_raw = df_pub_raw[df_pub_raw["year"] <= max_pub_before - pub_trunc]
            if df_pub_raw.empty:
                return jsonify({"ok": False, "error": "No publication data left after truncation"}), 400
            max_pub_after = int(df_pub_raw["year"].max())
            eff_pub_trunc = max_pub_before - max_pub_after

            # ---- yearly counts + merge ----
            pat_counts = df_pat_raw.groupby('year').size().reset_index(name='patents').astype({'year': int})
            pub_counts = df_pub_raw.groupby('year').size().reset_index(name='publications').astype({'year': int})

            merged = (
                pd.merge(pat_counts, pub_counts, on='year', how='inner')
              .sort_values('year')
              .reset_index(drop=True)
            )
            if len(merged) < (n_steps + 1):
                return jsonify({
                    "ok": False,
                    "error": f"Not enough merged data after truncation; need ≥{n_steps+1} years, have {len(merged)}",
                    "years_available": merged["year"].tolist()
            }), 400

            years_all = merged["year"].tolist()
            y_min, y_max = int(merged["year"].min()), int(merged["year"].max())

            # ---- enhanced features ----
            df_enh, X_raw, years_vec, feature_cols = _build_enhanced(
                merged.loc[:, ["year", "patents", "publications"]]
            )
            if hasattr(x_scaler, "n_features_in_") and X_raw.shape[1] != x_scaler.n_features_in_:
                return jsonify({
                    "ok": False,
                    "error": f"Feature mismatch: X has {X_raw.shape[1]} cols, scaler expects {x_scaler.n_features_in_}"
                }), 500

            X_scaled = x_scaler.transform(X_raw)
            X_win, target_years = _make_windows(X_scaled, years_vec, n_steps)

            # ---- evaluation split ----
            if split_year is not None:
                split_year = max(y_min, min(int(split_year), y_max - 1))
                test_years_list = list(range(split_year + 1, y_max + 1))
                test_mask  = np.isin(target_years, test_years_list)
            else:
                if test_years > 0:
                    split_year = y_max - test_years
                    test_years_list = list(range(split_year + 1, y_max + 1))
                    test_mask  = np.isin(target_years, test_years_list)
                else:
                    test_years_list = []
                    test_mask  = np.zeros_like(target_years, dtype=bool)

            mspe = None
            test_points = []
            if np.any(test_mask):
                X_test = X_win[test_mask]
                yrs_t  = target_years[test_mask]
                y_scaled = model.predict(X_test, verbose=0).ravel()
                y_unscaled = y_scaler.inverse_transform(y_scaled.reshape(-1,1)).ravel()
                y_pred = _y_inv(y_unscaled)
                y_pred = np.maximum(y_pred, 0.0)
                y_true = merged.set_index("year").loc[yrs_t, "patents"].astype(float).values
                mspe = _mspe(y_true, y_pred)
                test_points = [
                    {"year": int(y), "actual": float(a), "yhat": float(p)}
                    for y, a, p in zip(yrs_t.tolist(), y_true.tolist(), y_pred.tolist())
                ]

        # ---- production forecast (recursive) ----
            last_year_hist = int(merged["year"].iloc[-1])
            if X_win.shape[0] == 0:
                return jsonify({"ok": False, "error": "Not enough rows to form a window."}), 400
            window = X_win[-1]  # (n_steps, n_features)
            hist_df = merged.loc[:, ["year", "patents", "publications"]].copy()
            last_pub_unscaled = float(hist_df["publications"].iloc[-1])

            future = []
            for i in range(horizon):
                X_input  = window.reshape(1, n_steps, window.shape[1])
                y_scaled = model.predict(X_input, verbose=0)[0][0]
                y_unscaled = float(y_scaler.inverse_transform([[y_scaled]])[0, 0])
                y_pred     = float(_y_inv(y_unscaled))
                y_pred     = max(0.0, y_pred)

                fut_year = last_year_hist + i + 1
                future.append({"year": fut_year, "predicted_patents": y_pred})

                hist_df = pd.concat([
                    hist_df,
                    pd.DataFrame([{"year": fut_year, "patents": y_pred, "publications": last_pub_unscaled}])
                ], ignore_index=True)

                df_e2, X_raw2, years2, _ = _build_enhanced(hist_df)
                X_scaled2 = x_scaler.transform(X_raw2)
                window = X_scaled2[-n_steps:, :]

        # ======= GROWTH RATES =======
            current_year = date.today().year

            pats_actual = (
                pat_counts.rename(columns={"patents":"patent_count"})
                          .loc[:, ["year","patent_count"]]
                          .sort_values("year")
                          .reset_index(drop=True)
            )

            # 1) Past + current (actual only)
            past_GR_a, past_s_a, past_e_a, past_lbl_a, curr_GR_a, curr_s_a, curr_e_a, curr_lbl_a = \
                compute_patent_growth_from_counts(pats_actual)

            # 2) Current recomputed with forecast for [current_year, +1, +2] if available
            pats_curr_hybrid = pats_actual.copy()
            future_map = {int(r["year"]): float(r["predicted_patents"]) for r in future}
            for y in [current_year, current_year+1, current_year+2]:
                if y in future_map:
                    pats_curr_hybrid = _upsert_counts(pats_curr_hybrid, y, future_map[y])
            _, _, _, _, curr_GR_h, curr_s_h, curr_e_h, curr_lbl_h = \
                compute_patent_growth_from_counts(pats_curr_hybrid)

            # 3) Evaluation: past growth with model values (replace test years by predictions)
            eval_growth = None
            if len(test_points) > 0:
                pats_eval = pats_actual.copy()
                for tp in test_points:
                    pats_eval = _upsert_counts(pats_eval, int(tp["year"]), float(tp["yhat"]))
                past_GR_m, past_s_m, past_e_m, past_lbl_m, _, _, _, _ = \
                    compute_patent_growth_from_counts(pats_eval)

                eval_growth = {
                    "past_actual_percent": past_GR_a,
                    "past_model_percent": past_GR_m,
                    "delta_percent": (past_GR_m - past_GR_a) if np.isfinite(past_GR_a) and np.isfinite(past_GR_m) else None,
                    "window": {"start_year": past_s_m, "end_year": past_e_m},
                    "labels": {"actual": past_lbl_a, "model": past_lbl_m}
                }
            

            return jsonify({
                "ok": True,
                "params": {
                    "horizon": horizon,
                    "split_year": int(split_year) if split_year is not None else None,
                    "test_years": test_years if split_year is None else None
                },
                "applied_truncation": {
                    "pat_trunc_requested": pat_trunc,
                    "pub_trunc_requested": pub_trunc,
                    "pat_trunc_effective": eff_pat_trunc,
                    "pub_trunc_effective": eff_pub_trunc
                },
                "history": {
                    "years_used": years_all,
                    "last_history_year": y_max
                },
                "evaluation": {
                    "test_years": test_years_list if split_year is not None or test_years > 0 else [],
                    "mspe_percent": mspe,
                    "points": test_points,
                    "past_growth_comparison": eval_growth
                },
                "growth": {
                    "past_actual": {
                        "percent": past_GR_a,
                        "label": past_lbl_a,
                        "window": {"start_year": past_s_a, "end_year": past_e_a}
                    },
                    "current_actual": {
                        "percent": curr_GR_a,
                        "label": curr_lbl_a,
                        "window": {"start_year": curr_s_a, "end_year": curr_e_a}
                    },
                    "current_with_forecast": {
                        "percent": curr_GR_h,
                        "label": curr_lbl_h,
                        "window": {"start_year": curr_s_h, "end_year": curr_e_h},
                        "note": "current window recomputed with forecast for [current_year, +1, +2] where available"
                    }
                },
                "forecast": future,
                "start_year": last_year_hist + 1,
                "forecast_horizon": horizon
            }), 200

        except Exception as e:
            return jsonify({"ok": False, "error": str(e)}), 500


        
        
        
        

    @app.route('/api/lstm_forecast_sereis', methods=['GET', 'POST'])
    def lstm_forecast_series():
        try:
            payload    = request.get_json(silent=True) or {}
            horizon    = _get_int(payload, "horizon",   5, 1)
            pat_trunc  = _get_int(payload, "pat_trunc", 0, 0)
            pub_trunc  = _get_int(payload, "pub_trunc", 0, 0)
            test_years = _get_int(payload, "test_years", 0, 0)
            split_year = payload.get("split_year", request.args.get("split_year"))
            split_year = int(split_year) if split_year not in (None, "", "null") else None

            # ---- model + scalers (MUST match training) ----
            model_path  = os.path.join("backend/outputs", "lstm_patent_forecaster_enhanced_series.keras")
            scaler_path = os.path.join("backend/outputs", "global_scalers.pkl")
            model   = load_model(model_path)
            scalers = joblib.load(scaler_path)
            x_scaler = scalers["x_scaler"] if isinstance(scalers, dict) else scalers.x_scaler
            y_scaler = scalers["y_scaler"] if isinstance(scalers, dict) else scalers.y_scaler

            # infer step length from saved model
            n_steps = int((model.input_shape[1] if isinstance(model.input_shape, (list, tuple)) else 5) or 5)

            # ---- fetch DB data ----
            with engine.connect() as conn:
                df_pat_raw = pd.read_sql(text("""
                    SELECT first_filing_year AS year
                    FROM raw_patents
                    WHERE first_filing_year IS NOT NULL
                """), conn)
            if df_pat_raw.empty:
                return jsonify({"ok": False, "error": "No patent data found"}), 404

            with engine.connect() as conn:
                df_pub_raw = pd.read_sql(text("""
                    SELECT year
                    FROM research_data3
                    WHERE year IS NOT NULL
                """), conn)
            if df_pub_raw.empty:
                return jsonify({"ok": False, "error": "No publication data found"}), 404

            # ---- truncation ----
            max_pat_before = int(df_pat_raw["year"].max())
            if pat_trunc > 0:
                df_pat_raw = df_pat_raw[df_pat_raw["year"] <= max_pat_before - pat_trunc]
            if df_pat_raw.empty:
                return jsonify({"ok": False, "error": "No patent data left after truncation"}), 400
            max_pat_after = int(df_pat_raw["year"].max())
            eff_pat_trunc = max_pat_before - max_pat_after

            max_pub_before = int(df_pub_raw["year"].max())
            if pub_trunc > 0:
                df_pub_raw = df_pub_raw[df_pub_raw["year"] <= max_pub_before - pub_trunc]
            if df_pub_raw.empty:
                return jsonify({"ok": False, "error": "No publication data left after truncation"}), 400
            max_pub_after = int(df_pub_raw["year"].max())
            eff_pub_trunc = max_pub_before - max_pub_after

            # ---- yearly counts + merge ----
            pat_counts = df_pat_raw.groupby('year').size().reset_index(name='patents').astype({'year': int})
            pub_counts = df_pub_raw.groupby('year').size().reset_index(name='publications').astype({'year': int})

            merged = (
                pd.merge(pat_counts, pub_counts, on='year', how='inner')
                .sort_values('year')
                .reset_index(drop=True)
            )
            if len(merged) < (n_steps + 1):
                return jsonify({
                    "ok": False,
                    "error": f"Not enough merged data after truncation; need ≥{n_steps+1} years, have {len(merged)}",
                    "years_available": merged["year"].tolist()
                }), 400

            years_all = merged["year"].tolist()
            y_min, y_max = int(merged["year"].min()), int(merged["year"].max())

            # ---- enhanced features (same as training) ----
            df_enh, X_raw, years_vec, feature_cols = _build_enhanced(
                merged.loc[:, ["year", "patents", "publications"]]
            )
            if hasattr(x_scaler, "n_features_in_") and X_raw.shape[1] != x_scaler.n_features_in_:
                return jsonify({
                    "ok": False,
                    "error": f"Feature mismatch: X has {X_raw.shape[1]} cols, scaler expects {x_scaler.n_features_in_}"
                }), 500

            X_scaled = x_scaler.transform(X_raw)
            X_win, target_years = _make_windows(X_scaled, years_vec, n_steps)

            # ---- evaluation split ----
            if split_year is not None:
                split_year = max(y_min, min(int(split_year), y_max - 1))
                test_years_list = list(range(split_year + 1, y_max + 1))
                test_mask  = np.isin(target_years, test_years_list)
            else:
                if test_years > 0:
                    split_year = y_max - test_years
                    test_years_list = list(range(split_year + 1, y_max + 1))
                    test_mask  = np.isin(target_years, test_years_list)
                else:
                    test_years_list = []
                    test_mask  = np.zeros_like(target_years, dtype=bool)

            mspe = None
            test_points = []
            if np.any(test_mask):
                X_test = X_win[test_mask]
                yrs_t  = target_years[test_mask]
                y_scaled = model.predict(X_test, verbose=0).ravel()
                y_unscaled = y_scaler.inverse_transform(y_scaled.reshape(-1,1)).ravel()
                y_pred = _y_inv(y_unscaled)
                y_pred = np.maximum(y_pred, 0.0)
                y_true = merged.set_index("year").loc[yrs_t, "patents"].astype(float).values
                mspe = _mspe(y_true, y_pred)
                test_points = [
                    {"year": int(y), "actual": float(a), "yhat": float(p)}
                    for y, a, p in zip(yrs_t.tolist(), y_true.tolist(), y_pred.tolist())
                ]

            # ---- production forecast (recursive) ----
            last_year_hist = int(merged["year"].iloc[-1])
            if X_win.shape[0] == 0:
                return jsonify({"ok": False, "error": "Not enough rows to form a window."}), 400
            window = X_win[-1]  # (n_steps, n_features)
            hist_df = merged.loc[:, ["year", "patents", "publications"]].copy()
            last_pub_unscaled = float(hist_df["publications"].iloc[-1])

            future = []
            for i in range(horizon):
                X_input  = window.reshape(1, n_steps, window.shape[1])
                y_scaled = model.predict(X_input, verbose=0)[0][0]
                y_unscaled = float(y_scaler.inverse_transform([[y_scaled]])[0, 0])
                y_pred     = float(_y_inv(y_unscaled))
                y_pred     = max(0.0, y_pred)

                fut_year = last_year_hist + i + 1
                future.append({"year": fut_year, "predicted_patents": y_pred})

                hist_df = pd.concat([
                    hist_df,
                    pd.DataFrame([{"year": fut_year, "patents": y_pred, "publications": last_pub_unscaled}])
                ], ignore_index=True)

                df_e2, X_raw2, years2, _ = _build_enhanced(hist_df)
                X_scaled2 = x_scaler.transform(X_raw2)
                window = X_scaled2[-n_steps:, :]

            # ======= GROWTH RATES =======
            current_year = date.today().year

            pats_actual = (
                pat_counts.rename(columns={"patents":"patent_count"})
                          .loc[:, ["year","patent_count"]]
                          .sort_values("year")
                          .reset_index(drop=True)
            )

            # 1) Past + current (actual only)
            past_GR_a, past_s_a, past_e_a, past_lbl_a, curr_GR_a, curr_s_a, curr_e_a, curr_lbl_a = \
                compute_patent_growth_from_counts(pats_actual)

            # 2) Current recomputed with forecast for [current_year, +1, +2] if available
            pats_curr_hybrid = pats_actual.copy()
            future_map = {int(r["year"]): float(r["predicted_patents"]) for r in future}
            for y in [current_year, current_year+1, current_year+2]:
                if y in future_map:
                    pats_curr_hybrid = _upsert_counts(pats_curr_hybrid, y, future_map[y])
            _, _, _, _, curr_GR_h, curr_s_h, curr_e_h, curr_lbl_h = \
                compute_patent_growth_from_counts(pats_curr_hybrid)

            # 3) Evaluation: past growth with model values (replace test years by predictions)
            eval_growth = None
            if len(test_points) > 0:
                pats_eval = pats_actual.copy()
                for tp in test_points:
                    pats_eval = _upsert_counts(pats_eval, int(tp["year"]), float(tp["yhat"]))
                past_GR_m, past_s_m, past_e_m, past_lbl_m, _, _, _, _ = \
                    compute_patent_growth_from_counts(pats_eval)

                eval_growth = {
                    "past_actual_percent": past_GR_a,
                    "past_model_percent": past_GR_m,
                    "delta_percent": (past_GR_m - past_GR_a) if np.isfinite(past_GR_a) and np.isfinite(past_GR_m) else None,
                    "window": {"start_year": past_s_m, "end_year": past_e_m},
                    "labels": {"actual": past_lbl_a, "model": past_lbl_m}
                }

            # ======= HISTORY PAYLOADS (for plotting) =======
            patents_history = [
                {"year": int(r.year), "ds": _year_to_ds(int(r.year)), "count": int(r.patents)}
                for r in pat_counts.itertuples(index=False)
            ]
            publications_history = [
                {"year": int(r.year), "ds": _year_to_ds(int(r.year)), "count": int(r.publications)}
                for r in pub_counts.itertuples(index=False)
            ]
            merged_history = [
                {
                    "year": int(r.year),
                    "ds": _year_to_ds(int(r.year)),
                    "patents": int(r.patents),
                    "publications": int(r.publications)
                }
                for r in merged.itertuples(index=False)
            ]

            return jsonify({
                "ok": True,
                "params": {
                    "horizon": horizon,
                    "split_year": int(split_year) if split_year is not None else None,
                    "test_years": test_years if split_year is None else None
                },
                "applied_truncation": {
                    "pat_trunc_requested": pat_trunc,
                    "pub_trunc_requested": pub_trunc,
                    "pat_trunc_effective": eff_pat_trunc,
                    "pub_trunc_effective": eff_pub_trunc
                },
                "history": {
                    "patents": patents_history,
                    "publications": publications_history,
                    "merged": merged_history,
                    "years_used_for_model": years_all,
                    "last_history_year": y_max
                },
                "evaluation": {
                    "test_years": test_years_list if split_year is not None or test_years > 0 else [],
                    "mspe_percent": mspe,
                    "points": test_points,
                    "past_growth_comparison": eval_growth
                },
                "growth": {
                    "past_actual": {
                        "percent": past_GR_a,
                        "label": past_lbl_a,
                        "window": {"start_year": past_s_a, "end_year": past_e_a}
                    },
                    "current_actual": {
                        "percent": curr_GR_a,
                        "label": curr_lbl_a,
                        "window": {"start_year": curr_s_a, "end_year": curr_e_a}
                    },
                    "current_with_forecast": {
                        "percent": curr_GR_h,
                        "label": curr_lbl_h,
                        "window": {"start_year": curr_s_h, "end_year": curr_e_h},
                        "note": "current window recomputed with forecast for [current_year, +1, +2] where available"
                    }
                },
                "forecast": future,
                "start_year": last_year_hist + 1,
                "forecast_horizon": horizon
            }), 200

        except Exception as e:
            return jsonify({"ok": False, "error": str(e)}), 500

        
        
        
        
        
        
        
        
        
        
        
        
        
        
        
    
    
    
    @app.route('/api/patents/yearly_counts', methods=['GET'])
    def get_patent_yearly_counts():
        """
        Returns yearly counts of patents from raw_patents.
        Groups by first_filing_year.
        """
        try:
            query = text('''
                SELECT first_filing_year AS year, COUNT(*) AS count
                FROM raw_patents
                WHERE first_filing_year IS NOT NULL
                GROUP BY first_filing_year
                ORDER BY year
            ''')
            with engine.connect() as conn:
                result = conn.execute(query).fetchall()

            if not result:
                return jsonify({"message": "No patent data available"}), 404

            years = [int(r[0]) for r in result]
            counts = [r[1] for r in result]

            return jsonify({
                "labels": years,
                "datasets": [{
                    "label": "Patents per Year",
                    "data": counts,
                    
                }]
            }), 200
        except Exception as e:
            return jsonify({"error": str(e)}), 500


    @app.route('/api/publications/yearly_counts', methods=['GET'])
    def get_publication_yearly_counts():
        """
        Returns yearly counts of publications from research_data3.
        Groups by year column.
        """
        try:
            query = text('''
                SELECT year, COUNT(*) AS count
                FROM research_data3
                WHERE year IS NOT NULL
                GROUP BY year
                ORDER BY year
            ''')
            with engine.connect() as conn:
                result = conn.execute(query).fetchall()

            if not result:
                return jsonify({"message": "No publication data available"}), 404

            years = [int(r[0]) for r in result]
            counts = [r[1] for r in result]

            return jsonify({
                "labels": years,
                "datasets": [{
                    "label": "Publications per Year",
                    "data": counts,
                 
                }]
            }), 200
        except Exception as e:
            return jsonify({"error": str(e)}), 500

        

           
        
        
        
        
        
        

    @app.route('/api/report/generate-pptx', methods=['POST'])
    def generate_pptx():
        data = request.get_json()
        images = data.get('images', [])
        dimensions = data.get('dimensions', {})
    
        prs = Presentation()
    
        # Convert template dimensions to inches (1 inch = 72 points)
        template_width_in = dimensions.get('width', 0) / 72
        template_height_in = dimensions.get('height', 0) / 72
    
        # Set slide size based on template
        prs.slide_width = Inches(template_width_in)
        prs.slide_height = Inches(template_height_in)
    
        for img_obj in images:
            img_b64 = img_obj.get('data')
            img_width_px = img_obj.get('width', 0)
            img_height_px = img_obj.get('height', 0)
        
            # decode image
            image_stream = io.BytesIO(
                base64.b64decode(img_b64.split(',', 1)[1])
            )
        
            # Calculate aspect ratios
            img_aspect = img_width_px / img_height_px if img_height_px else 1
            slide_aspect = template_width_in / template_height_in
        
            # Calculate dimensions to maintain aspect ratio
            if img_aspect > slide_aspect:
                # Image is wider than slide - fit to width
                width = template_width_in
                height = template_width_in / img_aspect
                top = (template_height_in - height) / 2
                left = 0
            else:
                # Image is taller than slide - fit to height
                height = template_height_in
                width = template_height_in * img_aspect
                left = (template_width_in - width) / 2
                top = 0
            
            slide = prs.slides.add_slide(prs.slide_layouts[5])
        
            # Add image centered with correct aspect ratio
            slide.shapes.add_picture(
                image_stream,
                Inches(left),
                Inches(top),
                width=Inches(width),
                height=Inches(height)
            )


        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return send_file(
            output,
            as_attachment=True,
            download_name="report.pptx",
            mimetype=(
                "application/"
                "vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        )
        
        


    return app


if __name__ == '__main__':
    app = create_app() 
    

    try:
        with app.app_context():
            db.engine.connect()
            print('database conection successful!')
            db.create_all()
            print('database tables created successfully!')
    except Exception as e:
        print(f"Database connection failed: {e}")
        print(app.url_map)
        

    def find_free_port():
        with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
            s.bind(("", 0))
            return s.getsockname()[1]

    port = find_free_port()

    # Update or add PROT in .env
    env_path = os.path.join(os.path.dirname(__file__), '.env')
    try:
        with open(env_path, 'r') as f:
            lines = f.readlines()
        found = False
        for i, line in enumerate(lines):
            if re.match(r'^NEXT_PUBLIC_BACKEND_PORT\s*=.*', line):
                lines[i] = f'\n NEXT_PUBLIC_BACKEND_PORT={port}\n'
                found = True
                break
        if not found:
            lines.append(f'NEXT_PUBLIC_BACKEND_PORT={port}\n')
        with open(env_path, 'w') as f:
            f.writelines(lines)
        # Write port number to shared file for front and back access
        shared_path_file = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), 'shared_path.txt')
        with open(shared_path_file, 'w') as spf:
            spf.write(str(port))
        # Also copy port number to frontend/public/backend_port.txt for frontend access
        frontend_port_file = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'frontend', 'public', 'backend_port.txt')
        try:
            with open(frontend_port_file, 'w') as fpf:
                fpf.write(str(port))
        except Exception as fe:
            print(f"Failed to write frontend port file: {fe}")
    except Exception as e:
        print(f"Failed to update .env: {e}")

    server = ServerThread(app, port)
    server.start()

    def handle_exit(signum, frame):
        print("\nShutting down server...")
        server.shutdown()
        exit(0)

    signal.signal(signal.SIGINT, handle_exit)
    signal.signal(signal.SIGTERM, handle_exit)
